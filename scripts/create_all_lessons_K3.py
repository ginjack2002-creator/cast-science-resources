#!/usr/bin/env python3
"""ModelIt! Batch Generator - K-3 (Kindergarten through Grade 3)

Usage:
    python create_all_lessons_K3.py K     # Generate Kindergarten
    python create_all_lessons_K3.py 1     # Generate Grade 1
    python create_all_lessons_K3.py 2     # Generate Grade 2
    python create_all_lessons_K3.py 3     # Generate Grade 3
    python create_all_lessons_K3.py all   # Generate all K-3
"""
import os, sys, time, base64, requests
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
try:
    from pptx.enum.shapes import MSO_SHAPE
except:
    from pptx.util import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DI, Pt as DP, RGBColor as DR
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# API Key
import dotenv as _dotenv
_dotenv.load_dotenv(os.path.join(os.path.expanduser('~'), '.env'))
API_KEY = os.environ.get('OPENROUTER_API_KEY', '')
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# Brand colors (same as G04+)
NAVY=RGBColor(0x0D,0x1B,0x2A); BB=RGBColor(0x1A,0x47,0x80); MB=RGBColor(0x2E,0x86,0xAB)
LB=RGBColor(0x7E,0xC8,0xE3); SB=RGBColor(0x5D,0xB7,0xDE); OG=RGBColor(0xE6,0x7E,0x22)
DK=RGBColor(0x1A,0x1A,0x2E); WH=RGBColor(0xFF,0xFF,0xFF)
DN=DR(0x0D,0x1B,0x2A); DBB=DR(0x1A,0x47,0x80); DMB=DR(0x2E,0x86,0xAB)
total_cost=0.0

# Grade configuration
GRADE_CONFIG = {
    "K":  {"label": "Kindergarten", "age": "5-6 years old", "tier": 1,
           "dir": "grade-K",  "data": "lesson_data_GK",  "time": "25-35 minutes"},
    "1":  {"label": "1st Grade",    "age": "6-7 years old", "tier": 1,
           "dir": "grade-01", "data": "lesson_data_G01", "time": "25-35 minutes"},
    "2":  {"label": "2nd Grade",    "age": "7-8 years old", "tier": 2,
           "dir": "grade-02", "data": "lesson_data_G02", "time": "30-40 minutes"},
    "3":  {"label": "3rd Grade",    "age": "8-9 years old", "tier": 2,
           "dir": "grade-03", "data": "lesson_data_G03", "time": "35-45 minutes"},
}

# ── Image generation ─────────────────────────────────────────

def gen_img(scene, fn, img_dir, age_range, is_cartoon=False):
    global total_cost
    if is_cartoon:
        prompt = (f"Create a colorful, kid-friendly scientific diagram showing {scene}. "
                  "Style: clean educational illustration with bright colors, simple shapes, "
                  "labeled arrows, cartoon-style but scientifically accurate. White background, "
                  "bold outlines, easy to read labels. Suitable for elementary students.")
    else:
        prompt = (f"Create a photorealistic, cinematic, ultra-crisp image of {scene} "
                  "featuring a diverse, multicultural group of children with a wide variety "
                  "of ethnicities and backgrounds represented (White, Black, Latino, Asian, "
                  "Middle Eastern, mixed-race), in a bright, modern elementary classroom. "
                  f"Age-accurate for {age_range} (no one looks like an adult if they're a kid). "
                  "Style: warm, inviting, joyful. No stereotypes. Natural lighting, "
                  "professional editorial photo quality.")
    print(f"  Generating: {fn}...")
    try:
        r = requests.post("https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization":f"Bearer {API_KEY}","Content-Type":"application/json"},
            json={"model":"google/gemini-2.5-flash-image",
                  "messages":[{"role":"user","content":prompt}],
                  "modalities":["image","text"]},
            timeout=180)
        if r.status_code==200:
            data=r.json()
            imgs=data.get("choices",[{}])[0].get("message",{}).get("images",[])
            if imgs:
                img=imgs[0]; url=img.get("image_url",{}).get("url","") if isinstance(img,dict) else ""
                if url.startswith("data:image"):
                    fp=os.path.join(img_dir,fn)
                    with open(fp,"wb") as f: f.write(base64.b64decode(url.split(",")[1]))
                    total_cost+=data.get("usage",{}).get("cost",0)
                    print(f"    [OK] {fn}"); return fp
            print(f"    [!] No image in response"); return None
        else:
            print(f"    [X] {r.status_code}: {r.text[:100]}"); return None
    except Exception as e:
        print(f"    [X] {e}"); return None

def gen_all_imgs(c, img_dir, age_range):
    os.makedirs(img_dir, exist_ok=True); imgs={}
    for key,(fn,scene) in c["images"].items():
        fp=os.path.join(img_dir,fn)
        if os.path.exists(fp):
            print(f"    [SKIP] {fn}"); imgs[key]=fp
        else:
            is_cartoon = (key == "modeling")
            imgs[key]=gen_img(scene,fn,img_dir,age_range,is_cartoon); time.sleep(3)
    return imgs

# ── Shared utilities ──────────────────────────────────────────

def get_dims(path):
    try:
        with Image.open(path) as img: return img.width,img.height
    except: return None,None

def scaled(path,mw,mh):
    w,h=get_dims(path)
    if not w: return mw,mh
    r=w/h; sw=mw; sh=mw/r
    if sh>mh: sh=mh; sw=mh*r
    return sw,sh

def corners(slide,prs,pos="top"):
    if pos in["top","both"]:
        for shp,x,y,w,h,clr in[(MSO_SHAPE.RIGHT_TRIANGLE,-0.5,-0.5,4,1.5,NAVY),
            (MSO_SHAPE.PARALLELOGRAM,2.5,0,3,0.8,MB),(MSO_SHAPE.PARALLELOGRAM,3.5,0.5,2.5,0.4,LB)]:
            s=slide.shapes.add_shape(shp,Inches(x),Inches(y),Inches(w),Inches(h))
            s.fill.solid();s.fill.fore_color.rgb=clr;s.line.fill.background()
    if pos in["bottom","both"]:
        for shp,x,y,w,h,clr in[(MSO_SHAPE.RIGHT_TRIANGLE,6.5,6,4,1.5,NAVY),
            (MSO_SHAPE.PARALLELOGRAM,4.5,6.7,3,0.8,MB),(MSO_SHAPE.PARALLELOGRAM,4,6.3,2.5,0.4,LB)]:
            s=slide.shapes.add_shape(shp,Inches(x),Inches(y),Inches(w),Inches(h))
            s.fill.solid();s.fill.fore_color.rgb=clr;s.line.fill.background()
            if shp==MSO_SHAPE.RIGHT_TRIANGLE: s.rotation=180

def logo(slide,x,y,w=3):
    bg=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(0.9))
    bg.fill.solid();bg.fill.fore_color.rgb=SB;bg.line.fill.background()
    ic=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.15),Inches(y+0.25),Inches(0.4),Inches(0.4))
    ic.fill.solid();ic.fill.fore_color.rgb=WH;ic.line.fill.background()
    tb=slide.shapes.add_textbox(Inches(x+0.6),Inches(y+0.15),Inches(w-0.8),Inches(0.6))
    p=tb.text_frame.paragraphs[0];p.text="ModelIt!";p.font.size=Pt(32)
    p.font.bold=True;p.font.color.rgb=NAVY;p.font.name="Comic Sans MS"

def put_img(slide,path,x,y,mw,mh,center=False):
    if path and os.path.exists(path):
        sw,sh=scaled(path,mw,mh)
        if center: x+=(mw-sw)/2;y+=(mh-sh)/2
        slide.shapes.add_picture(path,Inches(x),Inches(y),Inches(sw),Inches(sh));return True
    b=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(mw),Inches(mh))
    b.fill.solid();b.fill.fore_color.rgb=RGBColor(0xF0,0xF5,0xFA);b.line.color.rgb=MB;return False

def pagenum(slide,n,tot=7):
    tb=slide.shapes.add_textbox(Inches(9.2),Inches(7.1),Inches(0.6),Inches(0.3))
    p=tb.text_frame.paragraphs[0];p.text=f"{n}/{tot}"
    p.font.size=Pt(10);p.font.color.rgb=RGBColor(0x66,0x66,0x66);p.alignment=PP_ALIGN.RIGHT

def cell_shd(cell,hex_c):
    shd=OxmlElement("w:shd");shd.set(qn("w:fill"),hex_c)
    cell._tc.get_or_add_tcPr().append(shd)

def name_date(doc, wide_lines=False):
    t=doc.add_table(rows=1,cols=2);t.alignment=WD_TABLE_ALIGNMENT.CENTER
    sz = DP(16) if wide_lines else DP(11)
    t.rows[0].cells[0].text="Name: _______________________________"
    t.rows[0].cells[1].text="Date: _______________"
    for c2 in t.rows[0].cells:
        for r2 in c2.paragraphs[0].runs: r2.font.size=sz
    doc.add_paragraph()

def hdr(doc,text,lv=1):
    p=doc.add_paragraph();r=p.add_run(text);r.bold=True
    if lv==1: r.font.size=DP(24);r.font.color.rgb=DN
    elif lv==2: r.font.size=DP(18);r.font.color.rgb=DBB
    elif lv==3: r.font.size=DP(14);r.font.color.rgb=DMB
    return p

def rbox(doc,lines=3,prompt=""):
    if prompt: p=doc.add_paragraph();p.add_run(prompt).font.size=DP(11)
    t=doc.add_table(rows=1,cols=1);t.style="Table Grid"
    t.rows[0].cells[0].text="\n"*lines;doc.add_paragraph();return t

def wide_rbox(doc,lines=3,prompt=""):
    """Response box with wider-spaced lines for young writers."""
    if prompt:
        p=doc.add_paragraph();r=p.add_run(prompt);r.font.size=DP(14);r.bold=True
    t=doc.add_table(rows=1,cols=1);t.style="Table Grid"
    cell=t.rows[0].cells[0]
    for i in range(lines):
        pp=cell.add_paragraph()
        pp.paragraph_format.space_after=DP(6)
        pp.paragraph_format.line_spacing=DP(28)
        r=pp.add_run("_" * 65)
        r.font.size=DP(14);r.font.color.rgb=DR(0xCC,0xCC,0xCC)
    doc.add_paragraph();return t

def drawing_box(doc, lines=8, prompt=""):
    """Large drawing box for K-1 students."""
    if prompt:
        p=doc.add_paragraph();r=p.add_run(prompt);r.font.size=DP(14);r.bold=True
    t=doc.add_table(rows=1,cols=1);t.style="Table Grid"
    t.rows[0].cells[0].text="\n"*lines
    doc.add_paragraph();return t

def sentence_frame_box(doc, frame_text):
    """Dotted-midline writing area with sentence frame for K-1."""
    p=doc.add_paragraph()
    r=p.add_run(frame_text);r.font.size=DP(16);r.bold=True;r.font.color.rgb=DBB
    doc.add_paragraph()
    for i in range(3):
        p=doc.add_paragraph()
        p.paragraph_format.space_after=DP(4)
        r=p.add_run("_ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _")
        r.font.size=DP(16);r.font.color.rgb=DR(0xBB,0xBB,0xBB)
        p2=doc.add_paragraph()
        p2.paragraph_format.space_before=DP(0)
        r2=p2.add_run("- - - - - - - - - - - - - - - - - - - - - - - - - - - - - -")
        r2.font.size=DP(10);r2.font.color.rgb=DR(0xDD,0xDD,0xDD)

# ══════════════════════════════════════════════════════════════
#  TIER 1: K-1 POWERPOINT (7 SLIDES)
# ══════════════════════════════════════════════════════════════

def make_ppt_k1(c, imgs, out_dir, grade_label):
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    lid = c["id"]; title = c["title"]; sub = c["subtitle"]; ngss = c["ngss"]

    # ── Slide 1: Cover ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "both"); logo(sl, 3.5, 1.3)
    put_img(sl, imgs.get("cover"), 0.3, 4.5, 3.2, 2.5, center=True)

    tb = sl.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.6))
    p = tb.text_frame.paragraphs[0]; p.text = "Student Lesson"
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    tb2 = sl.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(1.2))
    tf = tb2.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = BB; p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph(); p.text = sub
    p.font.size = Pt(15); p.font.italic = True; p.font.color.rgb = DK; p.alignment = PP_ALIGN.CENTER

    badge = sl.shapes.add_textbox(Inches(6), Inches(5.5), Inches(3.5), Inches(0.8))
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = f"NGSS: {ngss}"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = MB; p.alignment = PP_ALIGN.RIGHT
    p = tf.add_paragraph(); p.text = grade_label
    p.font.size = Pt(12); p.font.color.rgb = DK; p.alignment = PP_ALIGN.RIGHT
    pagenum(sl, 1, 7)

    # ── Slide 2: What Will We Learn? + Big Question (combined) ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "What Will We Learn?"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    goals = sl.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.5), Inches(2.5))
    tf = goals.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Today We Will:"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = BB
    for obj in c["objectives"]:
        p = tf.add_paragraph(); p.text = f"  *  {obj}"
        p.font.size = Pt(16); p.font.color.rgb = DK; p.space_before = Pt(8)

    # Big Question on same slide
    qbox = sl.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(5.5), Inches(2.0))
    tf = qbox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Our Big Question:"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = OG
    p = tf.add_paragraph(); p.text = c["big_question"]
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = BB; p.space_before = Pt(8)

    put_img(sl, imgs.get("landscape"), 6.0, 2.5, 3.5, 4.0, center=True)
    pagenum(sl, 2, 7)

    # ── Slide 3: Let's Be Scientists! (simplified LEVER) ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "Let's Be Scientists!"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    steps = sl.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(5.2), Inches(4.5))
    tf = steps.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Scientists build MODELS to understand how things work!"
    p.font.size = Pt(16); p.font.italic = True; p.font.color.rgb = MB

    for icon, step, desc in [
        ("\U0001F50D", "1. FIND the Parts", "What are the pieces of our system?"),
        ("\U0001F517", "2. CONNECT the Parts", "How do the parts affect each other?"),
        ("\U0001F3AE", "3. TEST It!", "Run our model and see what happens!")]:
        p = tf.add_paragraph(); p.text = f"{icon}  {step}"
        p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = BB; p.space_before = Pt(16)
        p = tf.add_paragraph(); p.text = f"      {desc}"
        p.font.size = Pt(16); p.font.color.rgb = DK

    put_img(sl, imgs.get("discussion"), 5.8, 2.5, 3.8, 3.5, center=True)
    pagenum(sl, 3, 7)

    # ── Slide 4: What Are the Parts? (components + cartoon diagram) ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "What Are the Parts?"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(4.8), Inches(4.5))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Our System Has These Parts:"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = BB
    for name, desc, ext in c["components"]:
        label = "(Outside)" if ext else "(Inside)"
        p = tf.add_paragraph(); p.text = f"  *  {name} {label}"
        p.font.size = Pt(18); p.font.bold = True; p.space_before = Pt(10)
        p = tf.add_paragraph(); p.text = f"       {desc}"
        p.font.size = Pt(14); p.font.color.rgb = DK

    put_img(sl, imgs.get("modeling"), 5.6, 2.1, 4.0, 4.5, center=True)
    pagenum(sl, 4, 7)

    # ── Slide 5: How Are They Connected? ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "How Are They Connected?"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.5), Inches(3.5))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "When one thing changes, what happens to the other?"
    p.font.size = Pt(17); p.font.color.rgb = DK

    p = tf.add_paragraph(); p.text = "(+) Goes UP Together"
    p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = RGBColor(0x22, 0x8B, 0x22); p.space_before = Pt(14)
    p = tf.add_paragraph(); p.text = "     More of this = More of that!"
    p.font.size = Pt(15)

    p = tf.add_paragraph(); p.text = "(-) Goes UP and DOWN"
    p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = RGBColor(0xDC, 0x14, 0x3C); p.space_before = Pt(14)
    p = tf.add_paragraph(); p.text = "     More of this = Less of that!"
    p.font.size = Pt(15)

    think = sl.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(5.5), Inches(1.2))
    tf = think.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Think About It:"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BB
    p = tf.add_paragraph(); p.text = c["think_about_it"]
    p.font.size = Pt(15); p.font.italic = True; p.space_before = Pt(6)

    put_img(sl, imgs.get("modeling"), 5.8, 2.5, 3.8, 3.5, center=True)
    pagenum(sl, 5, 7)

    # ── Slide 6: What Did We Find Out? ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "What Did We Find Out?"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.5), Inches(4.0))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Our Model Showed Us:"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = BB
    for d in c["discoveries"]:
        p = tf.add_paragraph(); p.text = f"  *  {d}"
        p.font.size = Pt(15); p.font.color.rgb = DK; p.space_before = Pt(10)

    ans = sl.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(5.5), Inches(1.0))
    tf = ans.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"Answer: {c['answer']}"
    p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = MB

    put_img(sl, imgs.get("cover"), 6.2, 2.5, 3.3, 3.0, center=True)
    pagenum(sl, 6, 7)

    # ── Slide 7: STEM Challenge + Career ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = f"STEM Challenge: {c['stem_title']}"
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.2), Inches(3.8))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "YOUR MISSION!"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = OG
    p = tf.add_paragraph(); p.text = c["stem_mission"]
    p.font.size = Pt(15); p.space_before = Pt(10)

    p = tf.add_paragraph(); p.text = "\nThink About:"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BB; p.space_before = Pt(10)
    for q in c["stem_questions"]:
        p = tf.add_paragraph(); p.text = f"  *  {q}"
        p.font.size = Pt(14); p.space_before = Pt(4)

    put_img(sl, imgs.get("stem"), 5.8, 2.3, 3.8, 3.5, center=True)

    car = sl.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9.0), Inches(0.95))
    car.fill.solid(); car.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    tf = car.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]; p.text = "COOL CAREER:  "
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xA5, 0x00)
    run = p.add_run(); run.text = c["career"]
    run.font.size = Pt(11); run.font.bold = False; run.font.color.rgb = WH
    pagenum(sl, 7, 7)

    out = os.path.join(out_dir, f"{lid}-Student-Presentation.pptx")
    prs.save(out); print(f"[OK] PPT (K-1, 7 slides): {out}"); return out


# ══════════════════════════════════════════════════════════════
#  TIER 2: GRADES 2-3 POWERPOINT (8 SLIDES)
# ══════════════════════════════════════════════════════════════

def make_ppt_g23(c, imgs, out_dir, grade_label):
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    lid = c["id"]; title = c["title"]; sub = c["subtitle"]; ngss = c["ngss"]

    # ── Slide 1: Cover ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "both"); logo(sl, 3.5, 1.3)
    put_img(sl, imgs.get("cover"), 0.3, 4.5, 3.2, 2.5, center=True)

    tb = sl.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.6))
    p = tb.text_frame.paragraphs[0]; p.text = "Student Lesson"
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    tb2 = sl.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(1.2))
    tf = tb2.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = BB; p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph(); p.text = sub
    p.font.size = Pt(15); p.font.italic = True; p.font.color.rgb = DK; p.alignment = PP_ALIGN.CENTER

    badge = sl.shapes.add_textbox(Inches(6), Inches(5.5), Inches(3.5), Inches(0.8))
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = f"NGSS: {ngss}"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = MB; p.alignment = PP_ALIGN.RIGHT
    p = tf.add_paragraph(); p.text = grade_label
    p.font.size = Pt(12); p.font.color.rgb = DK; p.alignment = PP_ALIGN.RIGHT
    pagenum(sl, 1, 8)

    # ── Slide 2: What Will We Learn? ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "What Will We Learn Today?"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    goals = sl.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.5), Inches(4.5))
    tf = goals.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Learning Goals"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = BB
    for obj in c["objectives"]:
        p = tf.add_paragraph(); p.text = f"  *  {obj}"
        p.font.size = Pt(16); p.font.color.rgb = DK; p.space_before = Pt(8)

    vocab = sl.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(4.3), Inches(4.5))
    tf = vocab.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Key Vocabulary"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = BB
    for term, defn in c["vocabulary"]:
        p = tf.add_paragraph(); p.text = f"  {term}"
        p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = NAVY; p.space_before = Pt(8)
        p = tf.add_paragraph(); p.text = f"     {defn}"
        p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = DK
    pagenum(sl, 2, 8)

    # ── Slide 3: The Big Question ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "The Big Question"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    tb2 = sl.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.0))
    tf = tb2.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = c["big_question"]
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = BB; p.alignment = PP_ALIGN.CENTER

    ctx = sl.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(4.5), Inches(1.5))
    tf = ctx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"{sub}. Today we'll build a MODEL to discover the answer!"
    p.font.size = Pt(16); p.font.color.rgb = DK

    put_img(sl, imgs.get("landscape"), 5.3, 3.2, 4.2, 3.2, center=True)
    pagenum(sl, 3, 8)

    # ── Slide 4: Let's Build a Model! (simplified LEVER) ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "Let's Build a Model!"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    lever = sl.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.2), Inches(4.5))
    tf = lever.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "A model helps us understand how things work!"
    p.font.size = Pt(16); p.font.italic = True; p.font.color.rgb = MB

    for step, desc in [("1. FIND", "Identify the PARTS of the system"),
                       ("2. CONNECT", "Draw ARROWS showing how parts affect each other"),
                       ("3. BUILD", "Put your model together in ModelIt!"),
                       ("4. TEST", "Run SIMULATIONS to see what happens"),
                       ("5. IMPROVE", "Make your model better with new evidence")]:
        p = tf.add_paragraph(); p.text = step
        p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = BB; p.space_before = Pt(10)
        p = tf.add_paragraph(); p.text = f"     {desc}"
        p.font.size = Pt(14); p.font.color.rgb = DK

    put_img(sl, imgs.get("modeling"), 5.8, 2.5, 3.8, 3.5, center=True)
    pagenum(sl, 4, 8)

    # ── Slide 5: Activity 1: Sort the Components ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "Activity 1: Sort the Parts"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    inst = sl.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(4.8), Inches(1.0))
    tf = inst.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sort these parts into OUTSIDE (inputs) and INSIDE (changes inside the system):"
    p.font.size = Pt(15); p.font.color.rgb = DK

    ct = sl.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(4.8), Inches(2.5))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Your Parts:"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = BB
    for name, desc, ext in c["components"]:
        p = tf.add_paragraph(); p.text = f"     *  {name}"
        p.font.size = Pt(16); p.space_before = Pt(6)

    think = sl.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(4.8), Inches(1.0))
    tf = think.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Think: Which parts can we control? Which happen on their own?"
    p.font.size = Pt(14); p.font.italic = True; p.font.color.rgb = MB

    put_img(sl, imgs.get("discussion"), 5.6, 2.1, 4.0, 3.8, center=True)
    pagenum(sl, 5, 8)

    # ── Slide 6: Activity 2: Connect with Arrows ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "Activity 2: Connect with Arrows"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.5), Inches(3))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Draw arrows to show HOW parts affect each other:"
    p.font.size = Pt(17)

    p = tf.add_paragraph(); p.text = "(+) POSITIVE"
    p.font.size = Pt(16); p.font.bold = True
    p.font.color.rgb = RGBColor(0x22, 0x8B, 0x22); p.space_before = Pt(14)
    p = tf.add_paragraph(); p.text = "     When one goes UP, the other goes UP too"
    p.font.size = Pt(14)

    p = tf.add_paragraph(); p.text = "(-) NEGATIVE"
    p.font.size = Pt(16); p.font.bold = True
    p.font.color.rgb = RGBColor(0xDC, 0x14, 0x3C); p.space_before = Pt(14)
    p = tf.add_paragraph(); p.text = "     When one goes UP, the other goes DOWN"
    p.font.size = Pt(14)

    ex = sl.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(5.5), Inches(1.5))
    tf = ex.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Think About It:"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BB
    p = tf.add_paragraph(); p.text = c["think_about_it"]
    p.font.size = Pt(15); p.font.italic = True; p.space_before = Pt(6)

    put_img(sl, imgs.get("discussion"), 5.8, 2.5, 3.8, 3.2, center=True)
    pagenum(sl, 6, 8)

    # ── Slide 7: Activity 3: Run the Simulation ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "Activity 3: Run the Simulation!"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5), Inches(4))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Test these scenarios in ModelIt!"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = BB
    for name, instr in c["scenarios"]:
        p = tf.add_paragraph(); p.text = name
        p.font.size = Pt(16); p.font.bold = True; p.space_before = Pt(12)
        p = tf.add_paragraph(); p.text = f"     {instr}"
        p.font.size = Pt(14)
    p = tf.add_paragraph(); p.text = "\nWatch the graphs change!"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = MB; p.space_before = Pt(16)

    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(2.1), Inches(4.3), Inches(4.2))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFA); box.line.color.rgb = MB
    tb2 = sl.shapes.add_textbox(Inches(5.5), Inches(4.0), Inches(3.9), Inches(1.0))
    p = tb2.text_frame.paragraphs[0]; p.text = "[ModelIt Platform Screenshot]"
    p.font.size = Pt(11); p.font.italic = True; p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER
    pagenum(sl, 7, 8)

    # ── Slide 8: Discoveries + STEM Challenge (combined) ──
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "What Did We Discover?"
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.5), Inches(2.5))
    tf = ct.text_frame; tf.word_wrap = True
    for d in c["discoveries"]:
        p = tf.add_paragraph(); p.text = f"  *  {d}"
        p.font.size = Pt(14); p.font.color.rgb = DK; p.space_before = Pt(4)

    # STEM mini-section
    stem = sl.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(5.5), Inches(2.5))
    tf = stem.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"STEM Challenge: {c['stem_title']}"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = OG
    p = tf.add_paragraph(); p.text = c["stem_mission"]
    p.font.size = Pt(13); p.space_before = Pt(6)

    put_img(sl, imgs.get("stem"), 6.0, 2.3, 3.6, 4.5, center=True)

    car = sl.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.5))
    car.fill.solid(); car.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    tf = car.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    p = tf.paragraphs[0]; p.text = "CAREER:  "
    p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xA5, 0x00)
    run = p.add_run(); run.text = c["career"]
    run.font.size = Pt(10); run.font.bold = False; run.font.color.rgb = WH
    pagenum(sl, 8, 8)

    out = os.path.join(out_dir, f"{lid}-Student-Presentation.pptx")
    prs.save(out); print(f"[OK] PPT (G2-3, 8 slides): {out}"); return out


# ══════════════════════════════════════════════════════════════
#  TIER 1: K-1 ACTIVITY PACK (6 pages, drawing/coloring)
# ══════════════════════════════════════════════════════════════

def make_pack_k1(c, out_dir, grade_label):
    lid=c["id"];title=c["title"];sub=c["subtitle"];ngss=c["ngss"]
    out=os.path.join(out_dir,f"{lid}-Student-Activity-Pack.docx")
    doc=Document()
    for s in doc.sections:
        s.top_margin=DI(0.7);s.bottom_margin=DI(0.7);s.left_margin=DI(0.7);s.right_margin=DI(0.7)

    # ── Page 1: Cover ──
    doc.add_paragraph();doc.add_paragraph()
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run("ModelIt!");r.bold=True;r.font.size=DP(52);r.font.color.rgb=DN;r.font.name="Comic Sans MS"
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run("My Activity Book");r.font.size=DP(28);r.font.color.rgb=DBB
    doc.add_paragraph()
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(title);r.bold=True;r.font.size=DP(30);r.font.color.rgb=DN
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(sub);r.italic=True;r.font.size=DP(15)
    doc.add_paragraph();doc.add_paragraph();doc.add_paragraph()
    name_date(doc, wide_lines=True)
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(f"{grade_label}  |  {ngss}");r.font.size=DP(14)
    doc.add_page_break()

    # ── Page 2: What Do You Know? (draw/circle pre-assessment) ──
    hdr(doc, "What Do You Know?")
    name_date(doc, wide_lines=True)
    p=doc.add_paragraph()
    r=p.add_run("Draw or circle your answer!")
    r.font.size=DP(14);r.italic=True
    doc.add_paragraph()

    for q in c["pre_assessment"][:2]:
        p=doc.add_paragraph();r=p.add_run(q);r.bold=True;r.font.size=DP(14)
        drawing_box(doc, lines=6, prompt="Draw your answer here:")
    doc.add_page_break()

    # ── Page 3: The Parts of Our System (cut & sort) ──
    hdr(doc, "The Parts of Our System")
    name_date(doc, wide_lines=True)
    p=doc.add_paragraph()
    r=p.add_run("Cut out the parts below. Sort them into the boxes!")
    r.font.size=DP(14);r.bold=True
    doc.add_paragraph()

    # Component cards
    hdr(doc, "Our Parts:", lv=2)
    for name, desc, ext in c["components"]:
        p=doc.add_paragraph()
        r=p.add_run(f"  \u2B50  {name}")
        r.bold=True;r.font.size=DP(16)
        p2=doc.add_paragraph()
        r2=p2.add_run(f"       {desc}")
        r2.font.size=DP(12)
    doc.add_paragraph()

    # Sorting boxes
    st = doc.add_table(rows=2, cols=2); st.style = "Table Grid"
    st.rows[0].cells[0].text = "OUTSIDE Parts\n(Come from outside)"
    st.rows[0].cells[1].text = "INSIDE Parts\n(Change inside system)"
    for cell in st.rows[0].cells:
        cell.paragraphs[0].runs[0].bold = True
        cell.paragraphs[0].runs[0].font.size = DP(13)
        cell_shd(cell, "7EC8E3")
    for cell in st.rows[1].cells:
        cell.text = "\n\n\n\n\n"
    doc.add_page_break()

    # ── Page 4: My Model Coloring Page ──
    hdr(doc, "My Model Coloring Page")
    name_date(doc, wide_lines=True)

    p=doc.add_paragraph()
    r=p.add_run("Color the parts of the system and trace the arrows!")
    r.font.size=DP(14);r.bold=True;r.font.color.rgb=DBB
    doc.add_paragraph()

    # Model description
    p=doc.add_paragraph()
    r=p.add_run(c.get("coloring_description", "Color each part of the model. Follow the arrows to see how they connect!"))
    r.font.size=DP(12);r.italic=True
    doc.add_paragraph()

    # Large coloring area (represented as a big box)
    ct = doc.add_table(rows=1, cols=1); ct.style = "Table Grid"
    cell = ct.rows[0].cells[0]
    # Build a text representation of the model
    ext_parts = [n for n, d, e in c["components"] if e]
    int_parts = [n for n, d, e in c["components"] if not e]
    model_text = "\n\n"
    for name in ext_parts:
        model_text += f"    [ {name} ]  ------>  "
    for i, name in enumerate(int_parts):
        model_text += f"[ {name} ]"
        if i < len(int_parts) - 1:
            model_text += "  ------>  "
    model_text += "\n\n\n    (+) means MORE of this = MORE of that"
    model_text += "\n    (-) means MORE of this = LESS of that"
    model_text += "\n\n"
    cell.text = model_text

    doc.add_paragraph()

    # Sentence frame with dotted-midline writing
    hdr(doc, "Write About Your Model:", lv=2)
    sentence_frame_box(doc, c.get("sentence_frame", "When _______ goes up, _______ goes _______."))
    doc.add_page_break()

    # ── Page 5: What Happened? (circle/draw for scenarios) ──
    hdr(doc, "What Happened?")
    name_date(doc, wide_lines=True)
    p=doc.add_paragraph()
    r=p.add_run("Circle or draw what you saw!")
    r.font.size=DP(14);r.italic=True
    doc.add_paragraph()

    for sname, sinstr in c["scenarios"][:2]:
        hdr(doc, sname, lv=2)
        p=doc.add_paragraph();r=p.add_run(sinstr);r.font.size=DP(12)
        doc.add_paragraph()

        # Smiley/frowny options
        p=doc.add_paragraph()
        r=p.add_run("Did it go UP or DOWN?    Circle:     \u2B06\uFE0F  UP     \u2B07\uFE0F  DOWN")
        r.font.size=DP(14);r.bold=True
        doc.add_paragraph()

        drawing_box(doc, lines=4, prompt="Draw what you saw:")

    doc.add_page_break()

    # ── Page 6: My Design! (STEM drawing) ──
    hdr(doc, f"My Design: {c['stem_title']}")
    name_date(doc, wide_lines=True)

    p=doc.add_paragraph()
    r=p.add_run(c["stem_mission"]);r.font.size=DP(13)
    doc.add_paragraph()

    drawing_box(doc, lines=10, prompt="Draw your design here:")

    # "I made..." sentence starter
    hdr(doc, "Tell about your design:", lv=2)
    p=doc.add_paragraph()
    r=p.add_run("I made "); r.font.size=DP(16); r.bold=True
    doc.add_paragraph()
    for i in range(3):
        p=doc.add_paragraph()
        r=p.add_run("_ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _ _")
        r.font.size=DP(16);r.font.color.rgb=DR(0xBB,0xBB,0xBB)

    doc.save(out);print(f"[OK] Activity Pack (K-1, 6 pages): {out}");return out


# ══════════════════════════════════════════════════════════════
#  TIER 2: GRADES 2-3 ACTIVITY PACK (8 pages, some writing)
# ══════════════════════════════════════════════════════════════

def make_pack_g23(c, out_dir, grade_label):
    lid=c["id"];title=c["title"];sub=c["subtitle"];ngss=c["ngss"]
    out=os.path.join(out_dir,f"{lid}-Student-Activity-Pack.docx")
    doc=Document()
    for s in doc.sections:
        s.top_margin=DI(0.7);s.bottom_margin=DI(0.7);s.left_margin=DI(0.7);s.right_margin=DI(0.7)

    # ── Page 1: Cover ──
    doc.add_paragraph();doc.add_paragraph()
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run("ModelIt!");r.bold=True;r.font.size=DP(48);r.font.color.rgb=DN
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run("Student Activity Pack");r.font.size=DP(28);r.font.color.rgb=DBB
    doc.add_paragraph()
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(title);r.bold=True;r.font.size=DP(30);r.font.color.rgb=DN
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(sub);r.italic=True;r.font.size=DP(15)
    doc.add_paragraph();doc.add_paragraph();doc.add_paragraph()
    name_date(doc, wide_lines=True)
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(f"Grade Level: {grade_label}  |  Standard: {ngss}");r.font.size=DP(12)
    doc.add_page_break()

    # ── Page 2: What Do You Know? ──
    hdr(doc,"What Do You Know?");name_date(doc, wide_lines=True)
    p=doc.add_paragraph();p.add_run("Answer these questions BEFORE the lesson.").font.size=DP(12)
    doc.add_paragraph()
    for q in c["pre_assessment"]:
        p=doc.add_paragraph();r=p.add_run(q);r.bold=True;r.font.size=DP(13)
        wide_rbox(doc, 3)
    doc.add_page_break()

    # ── Page 3: Vocabulary Match ──
    hdr(doc,"Vocabulary Match");name_date(doc, wide_lines=True)
    p=doc.add_paragraph()
    r=p.add_run("Draw a line from each word to its meaning!")
    r.font.size=DP(13);r.italic=True
    doc.add_paragraph()

    # Two-column table: terms on left, shuffled definitions on right
    vt = doc.add_table(rows=len(c["vocabulary"])+1, cols=3)
    vt.style = "Table Grid"
    vt.rows[0].cells[0].text = "Word"
    vt.rows[0].cells[1].text = ""  # space for drawing lines
    vt.rows[0].cells[2].text = "What It Means"
    for cell in [vt.rows[0].cells[0], vt.rows[0].cells[2]]:
        cell.paragraphs[0].runs[0].bold = True
        cell.paragraphs[0].runs[0].font.size = DP(13)
        cell_shd(cell, "7EC8E3")

    # Terms in order, definitions reversed for matching
    defs_shuffled = list(reversed([d for _, d in c["vocabulary"]]))
    for i, (term, defn) in enumerate(c["vocabulary"]):
        vt.rows[i+1].cells[0].text = f"\n{term}\n"
        vt.rows[i+1].cells[0].paragraphs[0].runs[0].bold = True
        vt.rows[i+1].cells[0].paragraphs[0].runs[0].font.size = DP(13)
        vt.rows[i+1].cells[1].text = "\n\n"  # line-drawing space
        vt.rows[i+1].cells[2].text = f"\n{defs_shuffled[i]}\n"
        for para in vt.rows[i+1].cells[2].paragraphs:
            for r2 in para.runs: r2.font.size = DP(11)
    doc.add_page_break()

    # ── Page 4: Sort the Parts ──
    hdr(doc,"Sort the Parts");name_date(doc, wide_lines=True)
    p=doc.add_paragraph()
    r=p.add_run("Sort these parts into OUTSIDE and INSIDE the system:")
    r.font.size=DP(13)
    doc.add_paragraph()

    hdr(doc, "Our Parts:", lv=2)
    for name, desc, ext in c["components"]:
        p=doc.add_paragraph();r=p.add_run(f"  *  {name}: ");r.bold=True;r.font.size=DP(13)
        p.add_run(desc).font.size=DP(12)
    doc.add_paragraph()

    ct=doc.add_table(rows=2,cols=2);ct.style="Table Grid"
    for i,lbl in enumerate(["OUTSIDE Parts\n(Come from outside)","INSIDE Parts\n(Change inside system)"]):
        ct.rows[0].cells[i].text=lbl;ct.rows[0].cells[i].paragraphs[0].runs[0].bold=True
        ct.rows[0].cells[i].paragraphs[0].runs[0].font.size=DP(12)
        cell_shd(ct.rows[0].cells[i],"7EC8E3")
    for cell in ct.rows[1].cells: cell.text="\n\n\n\n"
    doc.add_paragraph()
    wide_rbox(doc, 3, "Why did you sort them this way?")
    doc.add_page_break()

    # ── Page 5: Draw Your Model ──
    hdr(doc,"Draw Your Model");name_date(doc, wide_lines=True)
    p=doc.add_paragraph()
    r=p.add_run("Draw the parts as boxes. Connect them with arrows labeled + or -.")
    r.font.size=DP(12)
    doc.add_paragraph()
    hdr(doc,"Relationship Key:",lv=3)
    for lbl,desc in[("(+) POSITIVE:","When one goes UP, the other goes UP too"),
                    ("(-) NEGATIVE:","When one goes UP, the other goes DOWN")]:
        p=doc.add_paragraph();r=p.add_run(lbl+" ");r.bold=True;r.font.size=DP(12)
        p.add_run(desc).font.size=DP(12)
    doc.add_paragraph()
    drawing_box(doc, lines=10, prompt="My Model:")
    doc.add_page_break()

    # ── Page 6: What Happened? (simulation observations) ──
    hdr(doc,"What Happened?");name_date(doc, wide_lines=True)
    p=doc.add_paragraph();p.add_run("Run each scenario and write what you saw.").font.size=DP(12)
    doc.add_paragraph()
    for sname, sinstr, spred in c["sim_scenarios"]:
        hdr(doc, sname, lv=2)
        p=doc.add_paragraph();r=p.add_run("Settings: ");r.bold=True;r.font.size=DP(12)
        p.add_run(sinstr).font.size=DP(12)
        wide_rbox(doc, 2, f"My Prediction: {spred}")
        wide_rbox(doc, 2, "What I Saw:")
    doc.add_page_break()

    # ── Page 7: Think About It (2 reflections) ──
    hdr(doc,"Think About It");name_date(doc, wide_lines=True)
    p=doc.add_paragraph();p.add_run("Answer these questions about what you learned.").font.size=DP(12)
    doc.add_paragraph()
    for q in c["reflections"][:2]:
        p=doc.add_paragraph();r=p.add_run(q);r.bold=True;r.font.size=DP(13)
        wide_rbox(doc, 3)
    doc.add_page_break()

    # ── Page 8: STEM Challenge ──
    hdr(doc,f"STEM Challenge: {c['stem_title']}");name_date(doc, wide_lines=True)
    hdr(doc,"YOUR MISSION",lv=2)
    p=doc.add_paragraph();p.add_run(c["stem_mission"]).font.size=DP(12);doc.add_paragraph()
    hdr(doc,"THE SCENARIO",lv=2)
    p=doc.add_paragraph();p.add_run(c["stem_scenario"]).font.size=DP(12);doc.add_paragraph()

    drawing_box(doc, lines=8, prompt="Draw your design:")

    hdr(doc,"Explain your design:",lv=2)
    wide_rbox(doc, 3)

    doc.save(out);print(f"[OK] Activity Pack (G2-3, 8 pages): {out}");return out


# ══════════════════════════════════════════════════════════════
#  SHARED: TEACHER'S GUIDE (all K-3 grades)
# ══════════════════════════════════════════════════════════════

def make_guide(c, out_dir, grade_label, age_range, tier, session_time):
    lid=c["id"];title=c["title"];sub=c["subtitle"];ngss=c["ngss"]
    out=os.path.join(out_dir,f"{lid}-Teachers-Guide.docx")
    doc=Document()
    for s in doc.sections:
        s.top_margin=DI(0.6);s.bottom_margin=DI(0.6);s.left_margin=DI(0.6);s.right_margin=DI(0.6)

    # Cover
    doc.add_paragraph()
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run("ModelIt! Teacher's Guide");r.bold=True;r.font.size=DP(36);r.font.color.rgb=DN
    doc.add_paragraph()
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(f"{lid}: {title}");r.bold=True;r.font.size=DP(22)
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(sub);r.font.size=DP(14);r.italic=True
    doc.add_paragraph()

    grouping = "Whole class + pairs" if tier == 1 else "Pairs or small groups of 3"
    materials = "Activity Books, crayons/markers, scissors, glue sticks" if tier == 1 else "Student devices, Activity Pack, Internet access, crayons/markers"

    qr=[[" Grade Level", grade_label],
        ["NGSS Standard",f"{ngss}: {c['ngss_desc']}"],
        ["Time Required",f"{session_time} (can extend with STEM challenge)"],
        ["Materials", materials],
        ["Prep Required","Print Activity Packs; review lesson; prepare STEM materials"],
        ["Grouping", grouping]]
    t=doc.add_table(rows=len(qr),cols=2);t.style="Table Grid"
    for i,(lbl,val) in enumerate(qr):
        t.rows[i].cells[0].text=lbl
        for para in t.rows[i].cells[0].paragraphs:
            for r2 in para.runs: r2.bold=True;r2.font.size=DP(10)
        cell_shd(t.rows[i].cells[0],"7EC8E3")
        t.rows[i].cells[1].text=val
        for para in t.rows[i].cells[1].paragraphs:
            for r2 in para.runs: r2.font.size=DP(10)
    doc.add_page_break()

    # Table of Contents
    hdr(doc,"Guide Contents")
    toc_items = ["1. NGSS Standards Unpacking & CAST Alignment",
                 "2. Background Content Knowledge for Teachers",
                 "3. The LEVER Framework",
                 "4. Slide-by-Slide Facilitation Guide",
                 "5. Answer Key with Expected Student Responses",
                 "6. STEM Challenge Teacher Guidance",
                 "7. Differentiation & Extensions",
                 "8. Common Misconceptions & How to Address Them"]
    if tier == 1:
        toc_items.insert(4, "   (includes Coloring Page & Sentence Frame guidance)")
    for item in toc_items:
        p=doc.add_paragraph();p.add_run(item).font.size=DP(12)
    doc.add_page_break()

    # 1. NGSS Standards
    hdr(doc,"1. NGSS Standards Unpacking & CAST Alignment")
    hdr(doc,f"Performance Expectation: {ngss}",lv=2)
    p=doc.add_paragraph();r=p.add_run(c['ngss_desc']);r.italic=True;r.font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"Three-Dimensional Learning",lv=2)
    for dim_name,dim_title,description in c["dimensions"]:
        p=doc.add_paragraph();r=p.add_run(dim_name+": ");r.bold=True;r.font.size=DP(11)
        r=p.add_run(dim_title);r.italic=True;r.font.size=DP(11)
        p=doc.add_paragraph();p.add_run(description).font.size=DP(11);doc.add_paragraph()

    # CAST items (tier 2 only has full CAST, tier 1 gets simplified)
    if tier == 2 and "cast_items" in c:
        hdr(doc,"CAST Testing Alignment",lv=2)
        p=doc.add_paragraph();p.add_run("Students should be able to:").font.size=DP(11)
        for item in c["cast_items"]:
            p=doc.add_paragraph(style="List Bullet");p.add_run(item).font.size=DP(11)
        doc.add_paragraph()
        if "cast_questions" in c:
            p=doc.add_paragraph();r=p.add_run("Sample CAST-Style Questions:");r.bold=True;r.font.size=DP(11)
            for qtype,qtext in c["cast_questions"]:
                p=doc.add_paragraph();r=p.add_run(qtype+" ");r.bold=True;r.font.size=DP(10)
                p.add_run(qtext).font.size=DP(10)
    doc.add_page_break()

    # 2. Background Content
    hdr(doc,"2. Background Content Knowledge for Teachers")
    hdr(doc,f"{title} — Science Background",lv=2)
    p=doc.add_paragraph();p.add_run(c["background_intro"]).font.size=DP(11);doc.add_paragraph()
    for sec_title,content in c["background_sections"]:
        p=doc.add_paragraph();r=p.add_run(sec_title+": ");r.bold=True;r.font.size=DP(11);r.font.color.rgb=DBB
        p=doc.add_paragraph();p.add_run(content).font.size=DP(11);doc.add_paragraph()
    doc.add_page_break()

    # 3. LEVER Framework
    hdr(doc,"3. The LEVER Framework")
    if tier == 1:
        p=doc.add_paragraph();p.add_run(
            "For K-1 students, the LEVER framework is simplified to 3 kid-friendly steps: "
            "FIND, CONNECT, and TEST. This preserves the scientific modeling process while "
            "being accessible to young learners."
        ).font.size=DP(11);doc.add_paragraph()
        for phase,desc,example in[
            ("FIND the Parts","Students identify the components of the system.",c["lever_L"]),
            ("CONNECT the Parts","Students discover how parts affect each other (+ or -).",c["lever_E"]),
            ("TEST It!","Students observe what happens when parts change.",c["lever_Ev"])]:
            p=doc.add_paragraph();r=p.add_run(phase);r.bold=True;r.font.size=DP(12);r.font.color.rgb=DBB
            p=doc.add_paragraph();p.add_run(desc).font.size=DP(11)
            p=doc.add_paragraph();r=p.add_run("In this lesson: ");r.bold=True;r.font.size=DP(10);r.font.color.rgb=DMB
            p.add_run(example).font.size=DP(10);doc.add_paragraph()
    else:
        p=doc.add_paragraph();p.add_run("LEVER guides students through authentic scientific modeling:").font.size=DP(11);doc.add_paragraph()
        for phase,desc,example in[
            ("L - LOCATE the System","Students identify the components of the system.",c["lever_L"]),
            ("E - ESTABLISH Relationships","Students determine + and - relationships between components.",c["lever_E"]),
            ("V - VISUALIZE & Model","Students build their model in ModelIt!.",c["lever_V"]),
            ("E - EVALUATE Outcomes","Students run simulations and observe results.",c["lever_Ev"]),
            ("R - REVISE & Extend","Students improve their model based on evidence.",c["lever_R"])]:
            p=doc.add_paragraph();r=p.add_run(phase);r.bold=True;r.font.size=DP(12);r.font.color.rgb=DBB
            p=doc.add_paragraph();p.add_run(desc).font.size=DP(11)
            p=doc.add_paragraph();r=p.add_run("In this lesson: ");r.bold=True;r.font.size=DP(10);r.font.color.rgb=DMB
            p.add_run(example).font.size=DP(10);doc.add_paragraph()
    doc.add_page_break()

    # 4. Slide-by-Slide Guide
    slide_count = 7 if tier == 1 else 8
    hdr(doc,f"4. Slide-by-Slide Facilitation Guide ({slide_count} Slides)")
    p=doc.add_paragraph();p.add_run("Shows what students see, what to say, and what to watch for.").font.size=DP(11);doc.add_paragraph()
    for slide in c["slides_guide"]:
        p=doc.add_paragraph();r=p.add_run(f"{slide['num']}: {slide['title']}");r.bold=True;r.font.size=DP(14);r.font.color.rgb=DBB
        st=doc.add_table(rows=4,cols=2);st.style="Table Grid"
        for ri,(lbl,val) in enumerate([("What Students See:",slide['visual']),
                                        ("What to Say:",slide['say']),
                                        ("What to Do/Watch For:",slide['do']),
                                        ("Approximate Time:",slide['time'])]):
            st.rows[ri].cells[0].text=lbl;cell_shd(st.rows[ri].cells[0],"E8F4F8")
            st.rows[ri].cells[0].paragraphs[0].runs[0].bold=True
            st.rows[ri].cells[0].paragraphs[0].runs[0].font.size=DP(9)
            st.rows[ri].cells[1].text=val
            for para in st.rows[ri].cells[1].paragraphs:
                for r2 in para.runs: r2.font.size=DP(9)
        doc.add_paragraph()

    # K-1: Coloring page guidance
    if tier == 1:
        doc.add_paragraph()
        hdr(doc,"Coloring Page & Sentence Frame Guidance",lv=2)
        p=doc.add_paragraph();p.add_run(
            "The coloring page shows a horizontal model with labeled components and arrows. "
            "Students color each component and trace the arrows. Below the model, students "
            "complete a sentence frame using the dotted-midline writing lines."
        ).font.size=DP(11)
        doc.add_paragraph()
        p=doc.add_paragraph();r=p.add_run("Sentence Frame: ");r.bold=True;r.font.size=DP(11)
        p.add_run(c.get("sentence_frame", "When _______ goes up, _______ goes _______.")).font.size=DP(11)
        doc.add_paragraph()
        p=doc.add_paragraph();r=p.add_run("Tips for K-1 writers:");r.bold=True;r.font.size=DP(11)
        for tip in ["Model the sentence frame on the board first",
                     "Accept phonetic spelling — focus on the science concept",
                     "Allow students to draw instead of write if needed",
                     "Pair students for peer support during writing"]:
            p=doc.add_paragraph(style="List Bullet");p.add_run(tip).font.size=DP(10)

    doc.add_page_break()

    # 5. Answer Key
    hdr(doc,"5. Answer Key with Expected Student Responses")
    hdr(doc,"Component Sorting",lv=2)
    ext_names=[n for n,d,e in c["components"] if e]
    int_names=[n for n,d,e in c["components"] if not e]
    sort_t=doc.add_table(rows=2,cols=2);sort_t.style="Table Grid"
    sort_t.rows[0].cells[0].text="OUTSIDE (External)"
    sort_t.rows[0].cells[1].text="INSIDE (Internal)"
    for cell in sort_t.rows[0].cells:
        cell.paragraphs[0].runs[0].bold=True;cell_shd(cell,"1A4780")
        cell.paragraphs[0].runs[0].font.color.rgb=DR(0xFF,0xFF,0xFF)
        cell.paragraphs[0].runs[0].font.size=DP(10)
    sort_t.rows[1].cells[0].text=", ".join(ext_names)
    sort_t.rows[1].cells[1].text=", ".join(int_names)
    for cell in sort_t.rows[1].cells:
        if cell.paragraphs[0].runs: cell.paragraphs[0].runs[0].font.size=DP(10)
    doc.add_paragraph()
    p=doc.add_paragraph();r=p.add_run("Expected reasoning: ");r.bold=True;r.font.size=DP(10)
    p.add_run(c["sort_reasoning"]).font.size=DP(10);doc.add_paragraph()

    hdr(doc,"Relationships",lv=2)
    for conn,sign,expl in c["relationships"]:
        p=doc.add_paragraph();r=p.add_run(conn+" = ");r.bold=True;r.font.size=DP(11)
        r=p.add_run(sign);r.bold=True;r.font.size=DP(11)
        r.font.color.rgb=DR(0x22,0x8B,0x22) if "+" in sign else DR(0xDC,0x14,0x3C)
        p=doc.add_paragraph();p.add_run(f'"{expl}"').font.size=DP(10)
    doc.add_paragraph()

    hdr(doc,"Simulation Observations",lv=2)
    for sc,ans in c["sim_answers"]:
        p=doc.add_paragraph();r=p.add_run(sc);r.bold=True;r.font.size=DP(11);r.font.color.rgb=DBB
        p=doc.add_paragraph();p.add_run(ans).font.size=DP(10);doc.add_paragraph()

    if tier == 2 and "reflection_exemplars" in c:
        hdr(doc,"Reflection Question Exemplars",lv=2)
        for q,ans in c["reflection_exemplars"]:
            p=doc.add_paragraph();r=p.add_run("Q: "+q);r.bold=True;r.font.size=DP(11)
            p=doc.add_paragraph();p.add_run(ans).font.size=DP(10);doc.add_paragraph()
    doc.add_page_break()

    # 6. STEM Challenge Guidance
    hdr(doc,"6. STEM Challenge: Teacher Guidance")
    hdr(doc,"How to Introduce",lv=2)
    p=doc.add_paragraph();p.add_run(c["stem_intro"]).font.size=DP(11);doc.add_paragraph()
    hdr(doc,"Key Scientific Concepts",lv=2)
    for concept,expl in c["stem_concepts"]:
        p=doc.add_paragraph();r=p.add_run(concept+" ");r.bold=True;r.font.size=DP(11)
        p.add_run(expl).font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"Evaluating Student Designs",lv=2)
    for lv2,criteria in c["stem_eval"]:
        p=doc.add_paragraph();r=p.add_run(lv2+" ");r.bold=True;r.font.size=DP(11)
        p.add_run(criteria).font.size=DP(11)
    doc.add_page_break()

    # 7. Differentiation
    hdr(doc,"7. Differentiation & Extensions")
    hdr(doc,"For Struggling Learners",lv=2)
    for s2 in c["support"]:
        p=doc.add_paragraph(style="List Bullet");p.add_run(s2).font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"For Advanced Learners",lv=2)
    for e2 in c["extensions"]:
        p=doc.add_paragraph(style="List Bullet");p.add_run(e2).font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"Cross-Curricular Connections",lv=2)
    for subj,conn in c["cross_curr"]:
        p=doc.add_paragraph();r=p.add_run(subj+" ");r.bold=True;r.font.size=DP(11)
        p.add_run(conn).font.size=DP(11)
    doc.add_page_break()

    # 8. Misconceptions
    hdr(doc,"8. Common Misconceptions & How to Address Them")
    for mc,reality,resp in c["misconceptions"]:
        p=doc.add_paragraph();r=p.add_run("Misconception: ");r.bold=True;r.font.size=DP(11);r.font.color.rgb=DR(0xDC,0x14,0x3C)
        p.add_run(mc).font.size=DP(11)
        p=doc.add_paragraph();r=p.add_run("Reality: ");r.bold=True;r.font.size=DP(10)
        p.add_run(reality).font.size=DP(10)
        p=doc.add_paragraph();r=p.add_run("How to respond: ");r.italic=True;r.font.size=DP(10)
        p.add_run(resp).font.size=DP(10);doc.add_paragraph()
    doc.save(out);print(f"[OK] Teacher's Guide: {out}");return out


# ══════════════════════════════════════════════════════════════
#  MAIN: Lesson generation orchestrator
# ══════════════════════════════════════════════════════════════

def generate_lesson(c, grade_key):
    cfg = GRADE_CONFIG[grade_key]
    lid = c["id"]
    base_dir = os.path.join(SCRIPT_DIR, '..', 'materials', cfg["dir"])
    out_dir = os.path.join(base_dir, lid)
    img_dir = os.path.join(out_dir, "images")
    os.makedirs(out_dir, exist_ok=True); os.makedirs(img_dir, exist_ok=True)

    print(f"\n{'='*60}\nGenerating: {lid} - {c['title']} ({cfg['label']})\n{'='*60}")

    imgs = gen_all_imgs(c, img_dir, cfg["age"])

    if cfg["tier"] == 1:
        make_ppt_k1(c, imgs, out_dir, cfg["label"])
        make_pack_k1(c, out_dir, cfg["label"])
    else:
        make_ppt_g23(c, imgs, out_dir, cfg["label"])
        make_pack_g23(c, out_dir, cfg["label"])

    make_guide(c, out_dir, cfg["label"], cfg["age"], cfg["tier"], cfg["time"])
    print(f"[DONE] {lid} complete!\n")


def load_lessons(grade_key):
    """Dynamically import lesson data for a grade."""
    cfg = GRADE_CONFIG[grade_key]
    mod_name = cfg["data"]
    import importlib
    mod = importlib.import_module(mod_name)
    return mod.ALL_LESSONS


def run_grade(grade_key):
    cfg = GRADE_CONFIG[grade_key]
    lessons = load_lessons(grade_key)

    print("="*60)
    print(f"ModelIt! Batch Generator - {cfg['label']} ({len(lessons)} lessons)")
    print("="*60)
    for lesson in lessons:
        print(f"  - {lesson['id']}: {lesson['title']}")
    print()

    start_time = time.time()
    for i, lesson in enumerate(lessons, 1):
        print(f"\n[{i}/{len(lessons)}] Processing {lesson['id']}...")
        generate_lesson(lesson, grade_key)

    elapsed = time.time() - start_time
    print("\n" + "="*60)
    print(f"BATCH COMPLETE: {cfg['label']}")
    print("="*60)
    print(f"Lessons generated: {len(lessons)}")
    print(f"Total time: {elapsed/60:.1f} minutes")
    print(f"Estimated image cost: ${total_cost:.2f}")
    print("="*60)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python create_all_lessons_K3.py [K|1|2|3|all]")
        sys.exit(1)

    grade = sys.argv[1].upper()
    if grade == "ALL":
        for g in ["K", "1", "2", "3"]:
            run_grade(g)
    elif grade in GRADE_CONFIG:
        run_grade(grade)
    else:
        print(f"Unknown grade: {grade}. Use K, 1, 2, 3, or all.")
        sys.exit(1)
