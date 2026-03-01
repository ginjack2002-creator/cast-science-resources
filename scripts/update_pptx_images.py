#!/usr/bin/env python3
"""
Update PowerPoint images with diverse representation.

Regenerates 3 student-facing images per lesson (landscape, modeling, discussion)
using the updated diversity prompt template, then replaces them in the PPTX files.

Usage:
    python update_pptx_images.py --grade 8           # Update Grade 8 only
    python update_pptx_images.py --grade 5            # Update Grade 5 only
    python update_pptx_images.py --grade 6 7 8        # Update multiple grades
    python update_pptx_images.py --all                # Update all grades (5-10)
    python update_pptx_images.py --grade 8 --dry-run  # Preview without generating
"""
import os, sys, json, time, base64, argparse, requests, io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

import dotenv
dotenv.load_dotenv(Path.home() / '.env')

GOOGLE_AI_STUDIO_API_KEY = os.environ.get('GOOGLE_AI_STUDIO_API_KEY', '')
MATERIALS_DIR = Path(__file__).parent.parent / 'materials'

total_cost = 0.0
total_generated = 0

# Age ranges per grade
AGE_RANGES = {
    4: "9-10 years old",
    5: "10-11 years old",
    6: "11-12 years old",
    7: "12-13 years old",
    8: "13-14 years old",
    9: "14-15 years old",
    10: "15-16 years old",
}

# Diversity rotation: cycles through leadership for each image position
# Each grade gets a different starting rotation to ensure variety
DIVERSITY_LEADS = {
    "landscape": [
        "a Latina girl with long dark wavy hair wearing a green hoodie in the foreground",
        "a Black boy with a neat fade and glasses in the foreground",
        "a Vietnamese American girl with straight black hair and a bright t-shirt in the foreground",
        "a White boy with sandy brown hair and freckles in the foreground",
        "an Indian American girl with dark braids wearing a denim jacket in the foreground",
        "a Filipino boy with short black hair wearing a red hoodie in the foreground",
        "a Black girl with natural coily hair in puffs wearing a purple top in the foreground",
        "a White girl with curly auburn hair and a flannel shirt in the foreground",
        "a Korean American boy with glasses wearing a blue t-shirt in the foreground",
        "a Latino boy with curly dark hair wearing a gray hoodie in the foreground",
    ],
    "modeling": [
        "an East Asian male student with a fade haircut leading the group at his laptop",
        "a Latina girl with curly hair pointing at a model on screen",
        "a Black girl with braids wearing glasses, focused on her laptop",
        "a White boy with shaggy brown hair collaborating with the group",
        "a South Asian boy with short dark hair explaining his model to others",
        "a Black boy with short locs, typing on his laptop with concentration",
        "a White girl with a ponytail, pointing at data on screen",
        "a Chinese American girl with short bob haircut, working on her model",
        "a Latino boy with a buzz cut, sketching a diagram on paper",
        "a Japanese American boy with glasses, leading a group discussion",
    ],
    "discussion": [
        "a Black female teacher with natural hair leading the discussion",
        "a White male teacher with glasses engaging students in debate",
        "a Latina female teacher with dark hair facilitating conversation",
        "an Asian American male teacher encouraging student participation",
        "a Black male teacher with a warm smile moderating discussion",
        "a White female teacher with short blonde hair asking probing questions",
        "an Indian American female teacher guiding student thinking",
        "a Latino male teacher with a beard leading collaborative analysis",
        "an East Asian female teacher with long black hair drawing connections",
        "a Black female teacher with locs inspiring critical thinking",
    ],
}

# Occasional diversity features (applied every ~4th lesson)
OCCASIONAL_FEATURES = {
    2: "One student in the group uses a wheelchair.",
    5: "One middle school girl in the group wears a hijab with everyday clothes.",
    8: "One student in the group uses a wheelchair.",
}


def build_prompt(scene_desc, lead_desc, age_range, lesson_idx=0, extras=""):
    """Build a diversity-focused image prompt."""
    occasional = OCCASIONAL_FEATURES.get(lesson_idx, "")

    return f"""Create a photorealistic, cinematic, ultra-crisp image of {scene_desc}.

Feature {lead_desc}, with a genuinely diverse group of students around them — a mix of Latin American, Black, Asian/AAPI, and White students with varied facial features, different hairstyles (straight, wavy, curly, coily, braids, locs), and different body builds. No two people should look alike. Everyone wears modern everyday clothing — hoodies, t-shirts, jeans, sneakers — NO cultural or traditional outfits.{' ' + occasional if occasional else ''}{' ' + extras if extras else ''}

Age-accurate for {age_range} — students look their actual age, not like adults.
Style: modern education / future-ready / aspirational, realistic.
Composition: clean framing, strong focal subject, natural body proportions, believable clothing textures.
Camera/lighting: soft natural window light + gentle fill, shallow depth of field, 35mm/50mm look, sharp eyes, realistic skin texture.
Quality: high-resolution, crisp edges, minimal noise, professional editorial photo."""


def generate_image(prompt, output_path):
    """Generate an image using Google AI Studio."""
    global total_cost, total_generated

    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-image:generateContent?key={GOOGLE_AI_STUDIO_API_KEY}"

    resp = requests.post(url, json={
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {
            "responseModalities": ["IMAGE", "TEXT"]
        }
    }, timeout=120)

    if resp.status_code != 200:
        print(f"      [ERROR] ({resp.status_code}): {resp.text[:200]}")
        return False

    data = resp.json()
    for part in data.get("candidates", [{}])[0].get("content", {}).get("parts", []):
        if "inlineData" in part:
            img_b64 = part["inlineData"]["data"]
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'wb') as f:
                f.write(base64.b64decode(img_b64))
            total_cost += 0.039
            total_generated += 1
            return True

    print(f"      [ERROR] No image data in response")
    return False


def replace_image_in_pptx(pptx_path, old_image_name, new_image_path):
    """Replace an image in a PowerPoint file."""
    prs = Presentation(str(pptx_path))
    replaced = False

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                # Check if this is the image we want to replace
                image = shape.image
                if old_image_name.lower() in (image.content_type or '').lower() or True:
                    # Get the current position and size
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    # Can't directly replace image content in python-pptx easily
                    # Instead, we need to access the internal rels
                    pass

    # Alternative approach: use the image blob replacement
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    img = shape.image
                    # Check filename in the relationship
                    rel = None
                    for r in slide.part.rels.values():
                        if hasattr(r, 'target_ref') and old_image_name in str(r.target_ref):
                            rel = r
                            break

                    if rel:
                        # Read new image
                        with open(new_image_path, 'rb') as f:
                            new_data = f.read()
                        # Replace the blob
                        rel.target_part._blob = new_data
                        replaced = True
                except Exception:
                    pass

    if replaced:
        prs.save(str(pptx_path))
    return replaced


def update_pptx_images_direct(pptx_path, image_updates):
    """Replace images in a PPTX by matching on image index per slide.

    image_updates: dict mapping slide_index -> list of (new_image_path,) tuples
    """
    prs = Presentation(str(pptx_path))
    replaced_count = 0

    # Find all picture shapes and their image parts
    image_parts = {}
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture type
                try:
                    # Get the image part through the relationship
                    image_rId = shape._element.blipFill.blip.get(
                        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                    )
                    if image_rId:
                        rel = slide.part.rels[image_rId]
                        image_parts.setdefault(slide_idx, []).append((shape, rel))
                except Exception:
                    pass

    # Replace images for specified slides
    for slide_idx, new_paths in image_updates.items():
        if slide_idx in image_parts:
            shapes_rels = image_parts[slide_idx]
            for i, new_path in enumerate(new_paths):
                if i < len(shapes_rels):
                    shape, rel = shapes_rels[i]
                    try:
                        with open(new_path, 'rb') as f:
                            new_data = f.read()
                        rel.target_part._blob = new_data
                        replaced_count += 1
                    except Exception as e:
                        print(f"      [WARN] Could not replace image: {e}")

    if replaced_count > 0:
        prs.save(str(pptx_path))
    return replaced_count


def load_lesson_data(grade, level=None):
    """Load lesson data for a grade (and optional level for grades 9-10)."""
    scripts_dir = Path(__file__).parent

    if grade == 5:
        data_file = scripts_dir / "lesson_data_L02_L10.py"
    elif grade in (9, 10) and level:
        data_file = scripts_dir / f"lesson_data_G{grade:02d}_L{level}.py"
    else:
        data_file = scripts_dir / f"lesson_data_G{grade:02d}.py"

    if not data_file.exists():
        return {}

    ns = {}
    exec(data_file.read_text(encoding='utf-8'), ns)

    # Extract L01-L10 variables into a dict
    lessons = {}
    for i in range(1, 11):
        key = f"L{i:02d}"
        if key in ns:
            lessons[key] = ns[key]
    return lessons


def process_lesson_dirs(lesson_dirs, lessons, age_range, dry_run=False):
    """Process a list of lesson directories with given lesson data."""
    for lesson_idx, lesson_dir in enumerate(lesson_dirs):
        lesson_id = lesson_dir.name  # e.g., "G08-L01" or "G09L1-L01"
        pptx_files = list(lesson_dir.glob("*Student-Presentation.pptx"))
        if not pptx_files:
            print(f"  [SKIP] {lesson_id}: No PPTX found")
            continue

        pptx_path = pptx_files[0]
        images_dir = lesson_dir / "images"
        if not images_dir.exists():
            print(f"  [SKIP] {lesson_id}: No images directory")
            continue

        print(f"\n  {lesson_id}:")

        # Get lesson-specific image descriptions from lesson data
        lesson_key = f"L{lesson_idx+1:02d}"
        lesson_info = lessons.get(lesson_key, {})
        images_info = lesson_info.get("images", {})

        # Generate 3 new images: landscape, modeling, discussion
        image_types = ["landscape", "modeling", "discussion"]
        new_image_paths = {}

        for img_type in image_types:
            # Get scene description from lesson data
            img_data = images_info.get(img_type)
            if img_data:
                filename, scene_desc = img_data[0], img_data[1]
            else:
                # Fallback: find existing image file (try both prefix and suffix patterns)
                existing = list(images_dir.glob(f"{img_type}*"))
                if not existing:
                    existing = list(images_dir.glob(f"*-{img_type}*"))
                if not existing:
                    existing = list(images_dir.glob(f"*{img_type}*"))
                if existing:
                    filename = existing[0].name
                    scene_desc = f"students in a modern classroom engaged in {img_type} activity"
                else:
                    continue

            # Get diversity lead for this position
            leads = DIVERSITY_LEADS.get(img_type, DIVERSITY_LEADS["landscape"])
            lead = leads[lesson_idx % len(leads)]

            prompt = build_prompt(scene_desc, lead, age_range, lesson_idx)
            output_path = images_dir / filename

            if dry_run:
                print(f"    [DRY-RUN] Would regenerate: {filename}")
                print(f"              Lead: {lead[:60]}...")
                new_image_paths[img_type] = output_path
            else:
                print(f"    Generating: {filename}")
                if generate_image(prompt, output_path):
                    sz = output_path.stat().st_size / 1024
                    print(f"      [OK] {sz:.0f} KB")
                    new_image_paths[img_type] = output_path
                else:
                    print(f"      [FAIL]")
                time.sleep(1.5)

        # Update PPTX with new images
        if new_image_paths and not dry_run:
            # Map image types to slide indices:
            # Slide 2 (index 2) = landscape (Big Question)
            # Slide 3 (index 3) = modeling (Today We Build)
            # Slide 5 (index 5) = discussion (Activity 2)
            slide_image_map = {}
            if "landscape" in new_image_paths:
                slide_image_map[2] = [new_image_paths["landscape"]]
            if "modeling" in new_image_paths:
                slide_image_map[3] = [new_image_paths["modeling"]]
            if "discussion" in new_image_paths:
                slide_image_map[5] = [new_image_paths["discussion"]]

            if slide_image_map:
                count = update_pptx_images_direct(pptx_path, slide_image_map)
                print(f"    [PPTX] Updated {count} images in {pptx_path.name}")


def process_grade(grade, dry_run=False):
    """Process all lessons for a grade."""
    print(f"\n{'='*50}")
    print(f"Grade {grade}")
    print(f"{'='*50}")

    age_range = AGE_RANGES.get(grade, f"{grade+4}-{grade+5} years old")

    grade_dir = MATERIALS_DIR / f"grade-{grade:02d}"
    if not grade_dir.exists():
        print(f"  [SKIP] Grade directory not found: {grade_dir}")
        return

    # Grades 9-10 have level-based subdirectories
    if grade in (9, 10):
        for level in (1, 2, 3):
            level_dir = grade_dir / f"level-{level}"
            if not level_dir.exists():
                continue
            print(f"\n  --- Level {level} ---")
            lessons = load_lesson_data(grade, level=level)
            lesson_dirs = sorted([d for d in level_dir.iterdir() if d.is_dir()])
            process_lesson_dirs(lesson_dirs, lessons, age_range, dry_run)
    else:
        lessons = load_lesson_data(grade)
        lesson_dirs = sorted([d for d in grade_dir.iterdir() if d.is_dir()])
        process_lesson_dirs(lesson_dirs, lessons, age_range, dry_run)


def main():
    parser = argparse.ArgumentParser(description="Update PPTX images with diverse representation")
    parser.add_argument('--grade', type=int, nargs='+', help='Grade(s) to update (5-10)')
    parser.add_argument('--all', action='store_true', help='Update all grades')
    parser.add_argument('--dry-run', action='store_true', help='Preview without generating')
    args = parser.parse_args()

    if not args.grade and not args.all:
        parser.print_help()
        sys.exit(1)

    grades = list(range(4, 11)) if args.all else args.grade

    for grade in grades:
        process_grade(grade, dry_run=args.dry_run)

    print(f"\n{'='*50}")
    print(f"COMPLETE")
    print(f"{'='*50}")
    print(f"  Images generated: {total_generated}")
    print(f"  Estimated cost: ~${total_cost:.2f}")


if __name__ == '__main__':
    main()
