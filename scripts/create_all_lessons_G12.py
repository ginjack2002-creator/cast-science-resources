#!/usr/bin/env python3
"""
ModelIt! Batch Generator - Grade 12 (Levels 1, 2, 3)

Usage:
    python create_all_lessons_G12.py 1      # Generate Level 1 only
    python create_all_lessons_G12.py 2      # Generate Level 2 only
    python create_all_lessons_G12.py 3      # Generate Level 3 only
    python create_all_lessons_G12.py 1 2 3  # Generate all levels
    python create_all_lessons_G12.py        # Generate all levels (default)
"""
import os, sys, time, base64, requests
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
try:
    from pptx.enum.shapes import MSO_SHAPE
except:
    from pptx.util import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DI, Pt as DP, RGBColor as DR
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# API Key - read from ~/.env (never hardcode keys)
import dotenv as _dotenv
_dotenv.load_dotenv(os.path.join(os.path.expanduser('~'), '.env'))
API_KEY = os.environ.get('OPENROUTER_API_KEY', '')

NAVY=RGBColor(0x0D,0x1B,0x2A); BB=RGBColor(0x1A,0x47,0x80); MB=RGBColor(0x2E,0x86,0xAB)
LB=RGBColor(0x7E,0xC8,0xE3); SB=RGBColor(0x5D,0xB7,0xDE); OG=RGBColor(0xE6,0x7E,0x22)
DK=RGBColor(0x1A,0x1A,0x2E); WH=RGBColor(0xFF,0xFF,0xFF)
DN=DR(0x0D,0x1B,0x2A); DBB=DR(0x1A,0x47,0x80); DMB=DR(0x2E,0x86,0xAB)
total_cost=0.0

AGE_RANGE = "17-18 years old"

LEVEL_CONFIG = {
    1: {
        "label": "12th Grade \u2014 Level 1: Science, Society & You",
        "short": "Level 1: Science, Society & You",
        "base_dir": os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'materials', 'grade-12', 'level-1'),
        "import": "lesson_data_G12_L1",
    },
    2: {
        "label": "12th Grade \u2014 Level 2: Global Systems & Sustainability",
        "short": "Level 2: Global Systems & Sustainability",
        "base_dir": os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'materials', 'grade-12', 'level-2'),
        "import": "lesson_data_G12_L2",
    },
    3: {
        "label": "12th Grade \u2014 Level 3: Capstone Innovation Lab",
        "short": "Level 3: Capstone Innovation Lab",
        "base_dir": os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'materials', 'grade-12', 'level-3'),
        "import": "lesson_data_G12_L3",
    },
}

def gen_img(scene, fn, img_dir):
    global total_cost
    prompt = (f"Create a photorealistic, cinematic, ultra-crisp image of {scene} "
              "featuring a diverse, multicultural group with Black people centered "
              "(a mix of skin tones and ethnicities represented), age-accurate for "
              f"{AGE_RANGE} (no one looks like an adult if they're a kid). "
              "Style: modern education / future-ready / middle school coolness, "
              "confident, aspirational, realistic. Composition: clean framing, "
              "natural body proportions, realistic hair detail (coils, curls, locs, braids). "
              "Camera/lighting: soft natural window light, shallow depth of field, "
              "35mm/50mm look. Quality: high-resolution, professional editorial photo.")
    print(f"  Generating: {fn}...")
    try:
        r = requests.post("https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization":f"Bearer {API_KEY}","Content-Type":"application/json"},
            json={"model":"google/gemini-2.5-flash-image",
                  "messages":[{"role":"user","content":prompt}],
                  "modalities":["image","text"]},
            timeout=180)
        if r.status_code==200:
            data=r.json()
            imgs=data.get("choices",[{}])[0].get("message",{}).get("images",[])
            if imgs:
                img=imgs[0]; url=img.get("image_url",{}).get("url","") if isinstance(img,dict) else ""
                if url.startswith("data:image"):
                    fp=os.path.join(img_dir,fn)
                    with open(fp,"wb") as f: f.write(base64.b64decode(url.split(",")[1]))
                    total_cost+=data.get("usage",{}).get("cost",0)
                    print(f"    [OK] {fn}"); return fp
            print(f"    [!] No image in response"); return None
        else:
            print(f"    [X] {r.status_code}: {r.text[:100]}"); return None
    except Exception as e:
        print(f"    [X] {e}"); return None

def gen_all_imgs(c, img_dir):
    os.makedirs(img_dir,exist_ok=True); imgs={}
    for key,(fn,scene) in c["images"].items():
        fp=os.path.join(img_dir,fn)
        if os.path.exists(fp):
            print(f"    [SKIP] {fn}"); imgs[key]=fp
        else:
            imgs[key]=gen_img(scene,fn,img_dir); time.sleep(3)
    return imgs

def get_dims(path):
    try:
        with Image.open(path) as img: return img.width,img.height
    except: return None,None

def scaled(path,mw,mh):
    w,h=get_dims(path)
    if not w: return mw,mh
    r=w/h; sw=mw; sh=mw/r
    if sh>mh: sh=mh; sw=mh*r
    return sw,sh

def corners(slide,prs,pos="top"):
    if pos in["top","both"]:
        for shp,x,y,w,h,clr in[(MSO_SHAPE.RIGHT_TRIANGLE,-0.5,-0.5,4,1.5,NAVY),
            (MSO_SHAPE.PARALLELOGRAM,2.5,0,3,0.8,MB),(MSO_SHAPE.PARALLELOGRAM,3.5,0.5,2.5,0.4,LB)]:
            s=slide.shapes.add_shape(shp,Inches(x),Inches(y),Inches(w),Inches(h))
            s.fill.solid();s.fill.fore_color.rgb=clr;s.line.fill.background()
    if pos in["bottom","both"]:
        for shp,x,y,w,h,clr in[(MSO_SHAPE.RIGHT_TRIANGLE,6.5,6,4,1.5,NAVY),
            (MSO_SHAPE.PARALLELOGRAM,4.5,6.7,3,0.8,MB),(MSO_SHAPE.PARALLELOGRAM,4,6.3,2.5,0.4,LB)]:
            s=slide.shapes.add_shape(shp,Inches(x),Inches(y),Inches(w),Inches(h))
            s.fill.solid();s.fill.fore_color.rgb=clr;s.line.fill.background()
            if shp==MSO_SHAPE.RIGHT_TRIANGLE: s.rotation=180

def logo(slide,x,y,w=3):
    bg=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(0.9))
    bg.fill.solid();bg.fill.fore_color.rgb=SB;bg.line.fill.background()
    ic=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.15),Inches(y+0.25),Inches(0.4),Inches(0.4))
    ic.fill.solid();ic.fill.fore_color.rgb=WH;ic.line.fill.background()
    tb=slide.shapes.add_textbox(Inches(x+0.6),Inches(y+0.15),Inches(w-0.8),Inches(0.6))
    p=tb.text_frame.paragraphs[0];p.text="ModelIt!";p.font.size=Pt(32)
    p.font.bold=True;p.font.color.rgb=NAVY;p.font.name="Comic Sans MS"

def put_img(slide,path,x,y,mw,mh,center=False):
    if path and os.path.exists(path):
        sw,sh=scaled(path,mw,mh)
        if center: x+=(mw-sw)/2;y+=(mh-sh)/2
        slide.shapes.add_picture(path,Inches(x),Inches(y),Inches(sw),Inches(sh));return True
    b=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(mw),Inches(mh))
    b.fill.solid();b.fill.fore_color.rgb=RGBColor(0xF0,0xF5,0xFA);b.line.color.rgb=MB;return False

def pagenum(slide,n,tot=9):
    tb=slide.shapes.add_textbox(Inches(9.2),Inches(7.1),Inches(0.6),Inches(0.3))
    p=tb.text_frame.paragraphs[0];p.text=f"{n}/{tot}"
    p.font.size=Pt(10);p.font.color.rgb=RGBColor(0x66,0x66,0x66);p.alignment=PP_ALIGN.RIGHT

def cell_shd(cell,hex_c):
    shd=OxmlElement("w:shd");shd.set(qn("w:fill"),hex_c)
    cell._tc.get_or_add_tcPr().append(shd)

def name_date(doc):
    t=doc.add_table(rows=1,cols=2);t.alignment=WD_TABLE_ALIGNMENT.CENTER
    t.rows[0].cells[0].text="Name: _______________________________"
    t.rows[0].cells[1].text="Date: _______________"
    for c2 in t.rows[0].cells:
        for r2 in c2.paragraphs[0].runs: r2.font.size=DP(11)
    doc.add_paragraph()

def hdr(doc,text,lv=1):
    p=doc.add_paragraph();r=p.add_run(text);r.bold=True
    if lv==1: r.font.size=DP(24);r.font.color.rgb=DN
    elif lv==2: r.font.size=DP(18);r.font.color.rgb=DBB
    elif lv==3: r.font.size=DP(14);r.font.color.rgb=DMB
    return p

def rbox(doc,lines=3,prompt=""):
    if prompt: p=doc.add_paragraph();p.add_run(prompt).font.size=DP(11)
    t=doc.add_table(rows=1,cols=1);t.style="Table Grid"
    t.rows[0].cells[0].text="\n"*lines;doc.add_paragraph();return t

def make_ppt(c, imgs, out_dir, grade_label):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    lid = c["id"]; title = c["title"]; sub = c["subtitle"]; ngss = c["ngss"]

    # ========== SLIDE 1: COVER ==========
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "both")
    logo(sl, 3.5, 1.3)
    put_img(sl, imgs.get("cover"), 0.3, 4.5, 3.2, 2.5, center=True)

    tb = sl.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.6))
    p = tb.text_frame.paragraphs[0]; p.text = "Student Lesson"
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    tb2 = sl.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(1.2))
    tf = tb2.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = BB; p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph(); p.text = sub
    p.font.size = Pt(15); p.font.italic = True; p.font.color.rgb = DK; p.alignment = PP_ALIGN.CENTER

    badge = sl.shapes.add_textbox(Inches(6), Inches(5.5), Inches(3.5), Inches(0.8))
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = f"NGSS: {ngss}"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = MB; p.alignment = PP_ALIGN.RIGHT
    p = tf.add_paragraph(); p.text = grade_label
    p.font.size = Pt(12); p.font.color.rgb = DK; p.alignment = PP_ALIGN.RIGHT
    pagenum(sl, 1)

    # ========== SLIDE 2: LEARNING OBJECTIVES ==========
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "What You Will Learn Today"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    goals = sl.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.5), Inches(4.5))
    tf = goals.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Learning Goals"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = BB
    for obj in c["objectives"]:
        p = tf.add_paragraph(); p.text = f"  *  {obj}"
        p.font.size = Pt(16); p.font.color.rgb = DK; p.space_before = Pt(8)

    vocab = sl.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(4.3), Inches(4.5))
    tf = vocab.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Key Vocabulary"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = BB
    for term, defn in c["vocabulary"]:
        p = tf.add_paragraph(); p.text = f"  {term}"
        p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = NAVY; p.space_before = Pt(8)
        p = tf.add_paragraph(); p.text = f"     {defn}"
        p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = DK
    pagenum(sl, 2)

    # ========== SLIDE 3: THE BIG QUESTION ==========
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "The Big Question"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    tb2 = sl.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.0))
    tf = tb2.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = c["big_question"]
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = BB; p.alignment = PP_ALIGN.CENTER

    ctx = sl.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(4.5), Inches(1.5))
    tf = ctx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"{sub}. Today we'll build a MODEL to discover the answer!"
    p.font.size = Pt(16); p.font.color.rgb = DK
    put_img(sl, imgs.get("landscape"), 5.3, 3.2, 4.2, 3.2, center=True)
    pagenum(sl, 3)

    # ========== SLIDE 4: TODAY WE BUILD A MODEL ==========
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "Today We Will Build a Model!"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    lever = sl.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.2), Inches(4.5))
    tf = lever.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "A model helps us understand complex systems!"
    p.font.size = Pt(16); p.font.italic = True; p.font.color.rgb = MB
    for step, desc in [("1. LOCATE", "Identify the COMPONENTS (parts) of the system"),
                       ("2. ESTABLISH", "Connect them with RELATIONSHIPS (+ or -)"),
                       ("3. VISUALIZE", "Build your model in ModelIt!"),
                       ("4. EVALUATE", "Run SIMULATIONS to test scenarios"),
                       ("5. REVISE", "Improve your model based on evidence")]:
        p = tf.add_paragraph(); p.text = step
        p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = BB; p.space_before = Pt(10)
        p = tf.add_paragraph(); p.text = f"     {desc}"
        p.font.size = Pt(14); p.font.color.rgb = DK
    put_img(sl, imgs.get("modeling"), 5.8, 2.5, 3.8, 3.5, center=True)
    pagenum(sl, 4)

    # ========== SLIDE 5: ACTIVITY 1 - SORT COMPONENTS ==========
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "Activity 1: Sort the Components"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    inst = sl.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(4.8), Inches(1.0))
    tf = inst.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sort these components into EXTERNAL (inputs from outside) and INTERNAL (inside the system):"
    p.font.size = Pt(15); p.font.color.rgb = DK

    ct = sl.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(4.8), Inches(2.5))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Your Components:"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = BB
    for name, desc, ext in c["components"]:
        p = tf.add_paragraph(); p.text = f"     *  {name}"
        p.font.size = Pt(16); p.space_before = Pt(6)

    think = sl.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(4.8), Inches(1.0))
    tf = think.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Think: Which components can we control? Which happen on their own?"
    p.font.size = Pt(14); p.font.italic = True; p.font.color.rgb = MB
    put_img(sl, imgs.get("discussion"), 5.6, 2.1, 4.0, 3.8, center=True)
    pagenum(sl, 5)

    # ========== SLIDE 6: ACTIVITY 2 - CONNECT WITH ARROWS ==========
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "Activity 2: Connect with Arrows"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.5), Inches(3))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Draw arrows to show HOW components affect each other:"
    p.font.size = Pt(17)
    p = tf.add_paragraph(); p.text = "(+) POSITIVE Relationship"
    p.font.size = Pt(16); p.font.bold = True
    p.font.color.rgb = RGBColor(0x22, 0x8B, 0x22); p.space_before = Pt(14)
    p = tf.add_paragraph(); p.text = "     When one goes UP, the other goes UP too"
    p.font.size = Pt(14)
    p = tf.add_paragraph(); p.text = "(-) NEGATIVE Relationship"
    p.font.size = Pt(16); p.font.bold = True
    p.font.color.rgb = RGBColor(0xDC, 0x14, 0x3C); p.space_before = Pt(14)
    p = tf.add_paragraph(); p.text = "     When one goes UP, the other goes DOWN"
    p.font.size = Pt(14)

    ex = sl.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(5.5), Inches(1.5))
    tf = ex.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Think About It:"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BB
    p = tf.add_paragraph(); p.text = c["think_about_it"]
    p.font.size = Pt(15); p.font.italic = True; p.space_before = Pt(6)
    put_img(sl, imgs.get("discussion"), 5.8, 2.5, 3.8, 3.2, center=True)
    pagenum(sl, 6)

    # ========== SLIDE 7: ACTIVITY 3 - SIMULATE ==========
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "Activity 3: Run the Simulation!"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5), Inches(4))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Test these scenarios in ModelIt!"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = BB
    for name, instr in c["scenarios"]:
        p = tf.add_paragraph(); p.text = name
        p.font.size = Pt(16); p.font.bold = True; p.space_before = Pt(12)
        p = tf.add_paragraph(); p.text = f"     {instr}"
        p.font.size = Pt(14)
    p = tf.add_paragraph(); p.text = "\nWatch the activity graphs change!"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = MB; p.space_before = Pt(16)

    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(2.1), Inches(4.3), Inches(4.2))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFA); box.line.color.rgb = MB
    tb2 = sl.shapes.add_textbox(Inches(5.5), Inches(4.0), Inches(3.9), Inches(1.0))
    p = tb2.text_frame.paragraphs[0]; p.text = "[ModelIt Platform Screenshot - Simulation Results Graph]"
    p.font.size = Pt(11); p.font.italic = True; p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER
    pagenum(sl, 7)

    # ========== SLIDE 8: WHAT DID WE DISCOVER ==========
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = "What Did We Discover?"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.5), Inches(4.5))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Our Model Showed Us:"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = BB
    for d in c["discoveries"]:
        p = tf.add_paragraph(); p.text = f"  *  {d}"
        p.font.size = Pt(15); p.font.color.rgb = DK; p.space_before = Pt(10)

    ans = sl.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(5.5), Inches(1.2))
    tf = ans.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"Answer: {c['answer']}"
    p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = MB
    put_img(sl, imgs.get("cover"), 6.2, 2.5, 3.3, 3.0, center=True)
    pagenum(sl, 8)

    # ========== SLIDE 9: STEM CHALLENGE ==========
    sl = prs.slides.add_slide(blank)
    corners(sl, prs, "top")
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = f"STEM Challenge: {c['stem_title']}"
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER

    ct = sl.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.2), Inches(4.5))
    tf = ct.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "YOUR ENGINEERING MISSION"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = OG
    p = tf.add_paragraph(); p.text = c["stem_mission"]
    p.font.size = Pt(14); p.space_before = Pt(10)
    p = tf.add_paragraph(); p.text = "\nThe Challenge:"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BB; p.space_before = Pt(10)
    p = tf.add_paragraph(); p.text = c["stem_scenario"]
    p.font.size = Pt(14)
    p = tf.add_paragraph(); p.text = "\nThink Like an Engineer:"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BB; p.space_before = Pt(10)
    for q in c["stem_questions"]:
        p = tf.add_paragraph(); p.text = f"     *  {q}"
        p.font.size = Pt(13); p.space_before = Pt(4)
    put_img(sl, imgs.get("stem"), 5.8, 2.5, 3.8, 3.5, center=True)

    car = sl.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9.0), Inches(0.95))
    car.fill.solid(); car.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    tf = car.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]; p.text = "REAL CAREER CONNECTION:  "
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xA5, 0x00)
    run = p.add_run(); run.text = c["career"]
    run.font.size = Pt(11); run.font.bold = False; run.font.color.rgb = WH
    run2 = p.add_run()
    run2.text = " The skills you're using TODAY are the same ones they use on the job!"
    run2.font.size = Pt(11); run2.font.bold = False; run2.font.color.rgb = WH
    pagenum(sl, 9)

    out = os.path.join(out_dir, f"{lid}-Student-Presentation.pptx")
    prs.save(out); print(f"[OK] PPT: {out}"); return out

def make_pack(c, out_dir, grade_label):
    lid=c["id"];title=c["title"];sub=c["subtitle"];ngss=c["ngss"]
    out=os.path.join(out_dir,f"{lid}-Student-Activity-Pack.docx")
    doc=Document()
    for s in doc.sections:
        s.top_margin=DI(0.7);s.bottom_margin=DI(0.7);s.left_margin=DI(0.7);s.right_margin=DI(0.7)
    doc.add_paragraph();doc.add_paragraph()
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run("ModelIt!");r.bold=True;r.font.size=DP(48);r.font.color.rgb=DN
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run("Student Activity Pack");r.font.size=DP(28);r.font.color.rgb=DBB
    doc.add_paragraph()
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(title);r.bold=True;r.font.size=DP(30);r.font.color.rgb=DN
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(sub);r.italic=True;r.font.size=DP(15)
    doc.add_paragraph();doc.add_paragraph();doc.add_paragraph()
    name_date(doc)
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(f"Grade Level: {grade_label}  |  Standard: {ngss}");r.font.size=DP(12)
    doc.add_page_break()
    hdr(doc,"What's Inside");name_date(doc)
    for item,desc in[("1. Pre-Assessment","What do you already know?"),
                     ("2. Scoring Rubric","How your work will be assessed"),
                     ("3. Component Analysis","Sort and classify system parts"),
                     ("4. Model Recording Page","Sketch your model"),
                     ("5. Simulation Observations","Record what happens"),
                     ("6. Research & Extend","Add new components"),
                     ("7. Reflection Questions","What did you learn?"),
                     ("8. STEM Challenge",c["stem_title"])]:
        p=doc.add_paragraph();r=p.add_run(item);r.bold=True;r.font.size=DP(13)
        p.add_run(f"  -  {desc}").font.size=DP(11)
    doc.add_page_break()
    hdr(doc,"Pre-Assessment: What Do You Know?");name_date(doc)
    p=doc.add_paragraph();p.add_run("Answer these questions BEFORE the lesson begins.").font.size=DP(11)
    doc.add_paragraph()
    for q in c["pre_assessment"]:
        p=doc.add_paragraph();r=p.add_run(q);r.bold=True;r.font.size=DP(11);rbox(doc,3)
    doc.add_page_break()
    hdr(doc,"Scoring Rubric");name_date(doc)
    p=doc.add_paragraph();p.add_run("Use this rubric to check your own work!").font.size=DP(11);doc.add_paragraph()
    rdata=[["Skill","4 - Expert","3 - Proficient","2 - Developing","1 - Beginning"],
           ["Identifying\nComponents","All correct + explains role","Most correct","Some correct","Struggles to identify"],
           ["Establishing\nRelationships","All arrows correct + explains","Most correct","Some correct","Incorrect"],
           ["Running\nSimulations","Tests all; predicts outcomes","Runs all; records","Runs some with guidance","Needs help"],
           ["Explaining\nResults","Clear with evidence","Good explanation","Partial","Struggles"]]
    t=doc.add_table(rows=len(rdata),cols=5);t.style="Table Grid"
    for i,row in enumerate(rdata):
        for j,txt in enumerate(row):
            cell=t.rows[i].cells[j];cell.text=txt
            for para in cell.paragraphs:
                for r2 in para.runs: r2.font.size=DP(9)
            if i==0:
                cell.paragraphs[0].runs[0].bold=True;cell_shd(cell,"1A4780")
                cell.paragraphs[0].runs[0].font.color.rgb=DR(0xFF,0xFF,0xFF)
    doc.add_page_break()
    hdr(doc,"Component Analysis");name_date(doc)
    hdr(doc,"What is a Component?",lv=2)
    p=doc.add_paragraph();p.add_run("A ").font.size=DP(11)
    r=p.add_run("component");r.bold=True;r.font.size=DP(11)
    p.add_run(f" is a part of a system. In our {title} system, we have these components:").font.size=DP(11)
    doc.add_paragraph()
    for name,desc,ext in c["components"]:
        p=doc.add_paragraph();r=p.add_run(f"  *  {name}: ");r.bold=True;r.font.size=DP(11)
        p.add_run(desc).font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"Sorting Activity",lv=2)
    p=doc.add_paragraph();p.add_run("Sort into EXTERNAL (inputs from outside) and INTERNAL (processes inside the system):").font.size=DP(11)
    doc.add_paragraph()
    ct=doc.add_table(rows=2,cols=2);ct.style="Table Grid"
    for i,lbl in enumerate(["EXTERNAL Components\n(Come from outside)","INTERNAL Components\n(Happen inside system)"]):
        ct.rows[0].cells[i].text=lbl;ct.rows[0].cells[i].paragraphs[0].runs[0].bold=True
        cell_shd(ct.rows[0].cells[i],"7EC8E3")
    for cell in ct.rows[1].cells: cell.text="\n\n\n\n"
    doc.add_paragraph()
    p=doc.add_paragraph();r=p.add_run("Explain your reasoning:");r.bold=True;r.font.size=DP(11);rbox(doc,4)
    doc.add_page_break()
    hdr(doc,"Model Recording Page");name_date(doc)
    p=doc.add_paragraph();p.add_run("Sketch your model showing all components and their relationships.").font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"Relationship Key:",lv=3)
    for lbl,desc in[("(+) POSITIVE:","When one increases, the other ALSO increases"),
                    ("(-) NEGATIVE:","When one increases, the other DECREASES")]:
        p=doc.add_paragraph();r=p.add_run(lbl+" ");r.bold=True;r.font.size=DP(11)
        p.add_run(desc).font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"My Model Sketch:",lv=3)
    p=doc.add_paragraph();p.add_run("Draw components as boxes. Connect with arrows labeled + or -.").font.size=DP(10);p.italic=True
    sk=doc.add_table(rows=1,cols=1);sk.style="Table Grid";sk.rows[0].cells[0].text="\n"*11
    doc.add_paragraph()
    p=doc.add_paragraph();r=p.add_run("Explain ONE relationship in your model. Why does it work that way?")
    r.bold=True;r.font.size=DP(11);rbox(doc,4)
    doc.add_page_break()
    hdr(doc,"Simulation Observations");name_date(doc)
    p=doc.add_paragraph();p.add_run("Run each scenario in ModelIt! and record your observations.").font.size=DP(11);doc.add_paragraph()
    for sname,sinstr,spred in c["sim_scenarios"]:
        hdr(doc,sname,lv=2)
        p=doc.add_paragraph();r=p.add_run("Settings: ");r.bold=True;r.font.size=DP(11)
        p.add_run(sinstr).font.size=DP(11)
        p=doc.add_paragraph();r=p.add_run(spred);r.font.size=DP(11);r.italic=True;rbox(doc,2)
        p=doc.add_paragraph();r=p.add_run("Observation: What actually happened? Describe what you saw on the graph.")
        r.bold=True;r.font.size=DP(11);rbox(doc,3)
    doc.add_page_break()
    hdr(doc,"Research & Extend");name_date(doc)
    p=doc.add_paragraph();p.add_run("Real systems are more complex. EXTEND your model with new components!").font.size=DP(11);doc.add_paragraph()
    hdr(doc,"Possible New Components:",lv=2)
    for name,expl in c["extend_components"]:
        p=doc.add_paragraph();r=p.add_run(f"  *  {name}: ");r.bold=True;r.font.size=DP(11)
        p.add_run(expl).font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"My Extension:",lv=2)
    p=doc.add_paragraph();r=p.add_run("New component I want to add: ");r.bold=True;r.font.size=DP(11)
    p.add_run("________________________________").font.size=DP(11)
    p=doc.add_paragraph();r=p.add_run("How does it connect? (Draw and explain)");r.bold=True;r.font.size=DP(11)
    sk2=doc.add_table(rows=1,cols=1);sk2.style="Table Grid";sk2.rows[0].cells[0].text="\n"*5
    doc.add_paragraph()
    p=doc.add_paragraph();r=p.add_run("What evidence supports these connections?")
    r.bold=True;r.font.size=DP(11);rbox(doc,4)
    doc.add_page_break()
    hdr(doc,"Reflection Questions");name_date(doc)
    p=doc.add_paragraph();p.add_run("Answer after completing all activities.").font.size=DP(11);doc.add_paragraph()
    for q in c["reflections"]:
        p=doc.add_paragraph();r=p.add_run(q);r.bold=True;r.font.size=DP(11);rbox(doc,4)
    doc.add_page_break()
    hdr(doc,f"STEM CHALLENGE: {c['stem_title']}");name_date(doc)
    hdr(doc,"YOUR MISSION",lv=2)
    p=doc.add_paragraph();p.add_run(c["stem_mission"]).font.size=DP(11);doc.add_paragraph()
    hdr(doc,"THE SCENARIO",lv=2)
    p=doc.add_paragraph();p.add_run(c["stem_scenario"]).font.size=DP(11);doc.add_paragraph()
    hdr(doc,"ENGINEERING DESIGN THINKING",lv=2)
    for q in c["stem_design_qs"]:
        p=doc.add_paragraph(style="List Bullet");p.add_run(q).font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"MY DESIGN",lv=2)
    p=doc.add_paragraph();p.add_run("Draw your design here:").font.size=DP(10);p.italic=True
    db=doc.add_table(rows=1,cols=1);db.style="Table Grid";db.rows[0].cells[0].text="\n"*10
    doc.add_paragraph()
    hdr(doc,"EXPLAIN YOUR DESIGN",lv=2)
    p=doc.add_paragraph();r=p.add_run("Why did you design it this way? Use scientific reasoning!")
    r.bold=True;r.font.size=DP(11);rbox(doc,5)
    doc.save(out);print(f"[OK] Activity Pack: {out}");return out

def make_guide(c, out_dir, grade_label):
    lid=c["id"];title=c["title"];sub=c["subtitle"];ngss=c["ngss"]
    out=os.path.join(out_dir,f"{lid}-Teachers-Guide.docx")
    doc=Document()
    for s in doc.sections:
        s.top_margin=DI(0.6);s.bottom_margin=DI(0.6);s.left_margin=DI(0.6);s.right_margin=DI(0.6)
    doc.add_paragraph()
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run("ModelIt! Teacher's Guide");r.bold=True;r.font.size=DP(36);r.font.color.rgb=DN
    doc.add_paragraph()
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(f"{lid}: {title}");r.bold=True;r.font.size=DP(22)
    p=doc.add_paragraph();p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(sub);r.font.size=DP(14);r.italic=True
    doc.add_paragraph()
    qr=[[" Grade Level",grade_label],
        ["NGSS Standard",f"{ngss}: {c['ngss_desc']}"],
        ["Time Required","50-70 minutes (can extend to 90 min with STEM challenge)"],
        ["Materials","Student devices, Activity Pack, Internet access"],
        ["Prep Required","Load ModelIt! platform; print Activity Packs; review lesson"],
        ["Grouping","Pairs or small groups of 3 (recommended)"]]
    t=doc.add_table(rows=len(qr),cols=2);t.style="Table Grid"
    for i,(lbl,val) in enumerate(qr):
        t.rows[i].cells[0].text=lbl
        for para in t.rows[i].cells[0].paragraphs:
            for r2 in para.runs: r2.bold=True;r2.font.size=DP(10)
        cell_shd(t.rows[i].cells[0],"7EC8E3")
        t.rows[i].cells[1].text=val
        for para in t.rows[i].cells[1].paragraphs:
            for r2 in para.runs: r2.font.size=DP(10)
    doc.add_page_break()
    hdr(doc,"Guide Contents")
    for item in["1. NGSS Standards Unpacking & CAST Alignment",
                "2. Background Content Knowledge for Teachers",
                "3. The LEVER Framework","4. Slide-by-Slide Facilitation Guide",
                "5. Answer Key with Expected Student Responses",
                "6. STEM Challenge Teacher Guidance","7. Differentiation & Extensions",
                "8. Common Misconceptions & How to Address Them"]:
        p=doc.add_paragraph();p.add_run(item).font.size=DP(12)
    doc.add_page_break()
    hdr(doc,"1. NGSS Standards Unpacking & CAST Alignment")
    hdr(doc,f"Performance Expectation: {ngss}",lv=2)
    p=doc.add_paragraph();r=p.add_run(c['ngss_desc']);r.italic=True;r.font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"Three-Dimensional Learning",lv=2)
    for dim_name,dim_title,description in c["dimensions"]:
        p=doc.add_paragraph();r=p.add_run(dim_name+": ");r.bold=True;r.font.size=DP(11)
        r=p.add_run(dim_title);r.italic=True;r.font.size=DP(11)
        p=doc.add_paragraph();p.add_run(description).font.size=DP(11);doc.add_paragraph()
    hdr(doc,"CAST Testing Alignment",lv=2)
    p=doc.add_paragraph();p.add_run("The California Science Test (CAST) assesses this standard. Students should be able to:").font.size=DP(11)
    for item in c["cast_items"]:
        p=doc.add_paragraph(style="List Bullet");p.add_run(item).font.size=DP(11)
    doc.add_paragraph()
    p=doc.add_paragraph();r=p.add_run("Sample CAST-Style Questions:");r.bold=True;r.font.size=DP(11)
    for qtype,qtext in c["cast_questions"]:
        p=doc.add_paragraph();r=p.add_run(qtype+" ");r.bold=True;r.font.size=DP(10)
        p.add_run(qtext).font.size=DP(10)
    doc.add_page_break()
    hdr(doc,"2. Background Content Knowledge for Teachers")
    hdr(doc,f"{title} — Science Background",lv=2)
    p=doc.add_paragraph();p.add_run(c["background_intro"]).font.size=DP(11);doc.add_paragraph()
    for sec_title,content in c["background_sections"]:
        p=doc.add_paragraph();r=p.add_run(sec_title+": ");r.bold=True;r.font.size=DP(11);r.font.color.rgb=DBB
        p=doc.add_paragraph();p.add_run(content).font.size=DP(11);doc.add_paragraph()
    doc.add_page_break()
    hdr(doc,"3. The LEVER Framework")
    p=doc.add_paragraph();p.add_run("LEVER guides students through authentic scientific modeling:").font.size=DP(11);doc.add_paragraph()
    for phase,desc,example in[
        ("L - LOCATE the System","Students identify the components of the system.",c["lever_L"]),
        ("E - ESTABLISH Relationships","Students determine + and - relationships between components.",c["lever_E"]),
        ("V - VISUALIZE & Model","Students build their model in ModelIt!.",c["lever_V"]),
        ("E - EVALUATE Outcomes","Students run simulations and observe results.",c["lever_Ev"]),
        ("R - REVISE & Extend","Students improve their model based on evidence.",c["lever_R"])]:
        p=doc.add_paragraph();r=p.add_run(phase);r.bold=True;r.font.size=DP(12);r.font.color.rgb=DBB
        p=doc.add_paragraph();p.add_run(desc).font.size=DP(11)
        p=doc.add_paragraph();r=p.add_run("In this lesson: ");r.bold=True;r.font.size=DP(10);r.font.color.rgb=DMB
        p.add_run(example).font.size=DP(10);doc.add_paragraph()
    doc.add_page_break()
    hdr(doc,"4. Slide-by-Slide Facilitation Guide")
    p=doc.add_paragraph();p.add_run("Shows what students see, what to say, and what to watch for.").font.size=DP(11);doc.add_paragraph()
    for slide in c["slides_guide"]:
        p=doc.add_paragraph();r=p.add_run(f"{slide['num']}: {slide['title']}");r.bold=True;r.font.size=DP(14);r.font.color.rgb=DBB
        st=doc.add_table(rows=4,cols=2);st.style="Table Grid"
        for ri,(lbl,val) in enumerate([("What Students See:",slide['visual']),
                                        ("What to Say:",slide['say']),
                                        ("What to Do/Watch For:",slide['do']),
                                        ("Approximate Time:",slide['time'])]):
            st.rows[ri].cells[0].text=lbl;cell_shd(st.rows[ri].cells[0],"E8F4F8")
            st.rows[ri].cells[0].paragraphs[0].runs[0].bold=True
            st.rows[ri].cells[0].paragraphs[0].runs[0].font.size=DP(9)
            st.rows[ri].cells[1].text=val
            for para in st.rows[ri].cells[1].paragraphs:
                for r2 in para.runs: r2.font.size=DP(9)
        doc.add_paragraph()
    doc.add_page_break()
    hdr(doc,"5. Answer Key with Expected Student Responses")
    hdr(doc,"Component Sorting",lv=2)
    ext_names=[n for n,d,e in c["components"] if e]
    int_names=[n for n,d,e in c["components"] if not e]
    num_rows = max(len(ext_names), len(int_names)) + 1
    sort_t=doc.add_table(rows=num_rows,cols=2);sort_t.style="Table Grid"
    sort_t.rows[0].cells[0].text="EXTERNAL Components"
    sort_t.rows[0].cells[1].text="INTERNAL Components"
    for cell in sort_t.rows[0].cells:
        cell.paragraphs[0].runs[0].bold=True;cell_shd(cell,"1A4780")
        cell.paragraphs[0].runs[0].font.color.rgb=DR(0xFF,0xFF,0xFF)
        cell.paragraphs[0].runs[0].font.size=DP(10)
    for ri in range(1, num_rows):
        sort_t.rows[ri].cells[0].text = ext_names[ri-1] if ri-1 < len(ext_names) else ""
        sort_t.rows[ri].cells[1].text = int_names[ri-1] if ri-1 < len(int_names) else ""
        for cell in sort_t.rows[ri].cells:
            if cell.paragraphs[0].runs: cell.paragraphs[0].runs[0].font.size=DP(10)
    doc.add_paragraph()
    p=doc.add_paragraph();r=p.add_run("Expected student reasoning: ");r.bold=True;r.font.size=DP(10)
    p.add_run(c["sort_reasoning"]).font.size=DP(10);doc.add_paragraph()
    hdr(doc,"Relationships",lv=2)
    for conn,sign,expl in c["relationships"]:
        p=doc.add_paragraph();r=p.add_run(conn+" = ");r.bold=True;r.font.size=DP(11)
        r=p.add_run(sign);r.bold=True;r.font.size=DP(11)
        r.font.color.rgb=DR(0x22,0x8B,0x22) if "+" in sign else DR(0xDC,0x14,0x3C)
        p=doc.add_paragraph();r=p.add_run("Expected reasoning: ");r.italic=True;r.font.size=DP(10)
        p.add_run(f'"{expl}"').font.size=DP(10)
    doc.add_paragraph()
    hdr(doc,"Simulation Observations",lv=2)
    for sc,ans in c["sim_answers"]:
        p=doc.add_paragraph();r=p.add_run(sc);r.bold=True;r.font.size=DP(11);r.font.color.rgb=DBB
        p=doc.add_paragraph();p.add_run(ans).font.size=DP(10);doc.add_paragraph()
    hdr(doc,"Reflection Question Exemplars",lv=2)
    for q,ans in c["reflection_exemplars"]:
        p=doc.add_paragraph();r=p.add_run("Q: "+q);r.bold=True;r.font.size=DP(11)
        p=doc.add_paragraph();p.add_run(ans).font.size=DP(10);doc.add_paragraph()
    doc.add_page_break()
    hdr(doc,"6. STEM Challenge: Teacher Guidance")
    hdr(doc,"How to Introduce",lv=2)
    p=doc.add_paragraph();p.add_run(c["stem_intro"]).font.size=DP(11);doc.add_paragraph()
    hdr(doc,"Key Scientific Concepts",lv=2)
    for concept,expl in c["stem_concepts"]:
        p=doc.add_paragraph();r=p.add_run(concept+" ");r.bold=True;r.font.size=DP(11)
        p.add_run(expl).font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"Evaluating Student Designs",lv=2)
    for lv2,criteria in c["stem_eval"]:
        p=doc.add_paragraph();r=p.add_run(lv2+" ");r.bold=True;r.font.size=DP(11)
        p.add_run(criteria).font.size=DP(11)
    doc.add_page_break()
    hdr(doc,"7. Differentiation & Extensions")
    hdr(doc,"For Struggling Learners",lv=2)
    for s2 in c["support"]:
        p=doc.add_paragraph(style="List Bullet");p.add_run(s2).font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"For Advanced Learners",lv=2)
    for e2 in c["extensions"]:
        p=doc.add_paragraph(style="List Bullet");p.add_run(e2).font.size=DP(11)
    doc.add_paragraph()
    hdr(doc,"Cross-Curricular Connections",lv=2)
    for subj,conn in c["cross_curr"]:
        p=doc.add_paragraph();r=p.add_run(subj+" ");r.bold=True;r.font.size=DP(11)
        p.add_run(conn).font.size=DP(11)
    doc.add_page_break()
    hdr(doc,"8. Common Misconceptions & How to Address Them")
    for mc,reality,resp in c["misconceptions"]:
        p=doc.add_paragraph();r=p.add_run("Misconception: ");r.bold=True;r.font.size=DP(11);r.font.color.rgb=DR(0xDC,0x14,0x3C)
        p.add_run(mc).font.size=DP(11)
        p=doc.add_paragraph();r=p.add_run("Reality: ");r.bold=True;r.font.size=DP(10)
        p.add_run(reality).font.size=DP(10)
        p=doc.add_paragraph();r=p.add_run("How to respond: ");r.italic=True;r.font.size=DP(10)
        p.add_run(resp).font.size=DP(10);doc.add_paragraph()
    doc.save(out);print(f"[OK] Teacher's Guide: {out}");return out

def generate_lesson(c, base_dir, grade_label):
    lid=c["id"]
    out_dir=os.path.join(base_dir,lid)
    img_dir=os.path.join(out_dir,"images")
    os.makedirs(out_dir,exist_ok=True);os.makedirs(img_dir,exist_ok=True)
    print(f"\n{'='*60}\nGenerating: {lid} - {c['title']}\n{'='*60}")
    imgs=gen_all_imgs(c,img_dir)
    make_ppt(c,imgs,out_dir,grade_label)
    make_pack(c,out_dir,grade_label)
    make_guide(c,out_dir,grade_label)
    print(f"[DONE] {lid} complete!\n")

def run_level(level_num):
    global total_cost
    cfg = LEVEL_CONFIG[level_num]
    mod = __import__(cfg["import"])
    lessons = mod.ALL_LESSONS
    grade_label = cfg["label"]
    base_dir = cfg["base_dir"]

    print(f"\n{'='*60}")
    print(f"ModelIt! Batch Generator - Grade 12 {cfg['short']}")
    print(f"{'='*60}")
    print(f"\nGenerating {len(lessons)} lessons:")
    for lesson in lessons:
        print(f"  - {lesson['id']}: {lesson['title']}")
    print()

    start_time = time.time()
    level_cost_start = total_cost
    for i, lesson in enumerate(lessons, 1):
        print(f"\n[{i}/{len(lessons)}] Processing {lesson['id']}...")
        generate_lesson(lesson, base_dir, grade_label)

    elapsed = time.time() - start_time
    level_cost = total_cost - level_cost_start
    print(f"\n{'='*60}")
    print(f"Level {level_num} COMPLETE!")
    print(f"{'='*60}")
    print(f"Lessons generated: {len(lessons)}")
    print(f"Time: {elapsed/60:.1f} minutes")
    print(f"Level image cost: ${level_cost:.2f}")
    return len(lessons), elapsed, level_cost

if __name__ == "__main__":
    levels_to_run = [int(x) for x in sys.argv[1:]] if len(sys.argv) > 1 else [1, 2, 3]

    print("="*60)
    print("ModelIt! Grade 12 Generator")
    print(f"Levels to generate: {levels_to_run}")
    print("="*60)

    total_start = time.time()
    total_lessons = 0
    for level in levels_to_run:
        if level not in LEVEL_CONFIG:
            print(f"[!] Unknown level: {level}. Use 1, 2, or 3.")
            continue
        count, elapsed, cost = run_level(level)
        total_lessons += count

    total_elapsed = time.time() - total_start
    print("\n" + "="*60)
    print("ALL LEVELS COMPLETE!")
    print("="*60)
    print(f"Total lessons generated: {total_lessons}")
    print(f"Total time: {total_elapsed/60:.1f} minutes")
    print(f"Total estimated image cost: ${total_cost:.2f}")
    print("="*60)
