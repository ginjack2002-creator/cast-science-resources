#!/usr/bin/env python3
"""
ModelIt! Student Presentation Generator (Generic)
=================================================
Creates branded 9-slide PowerPoint presentations for any lesson
using data from lesson_data_*.py files.

Usage:
    python create_presentations.py --grade 04_L2 --all
    python create_presentations.py --grade 06_L2 --lesson 1
    python create_presentations.py --grade 08 --lesson 5
"""

import argparse
import os
import sys
import importlib.util

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# ============================================================
# Brand Colors
# ============================================================
PURPLE_DEEP = RGBColor(0x4A, 0x3A, 0x7A)
PURPLE_BRAND = RGBColor(0x6B, 0x5B, 0x95)
PURPLE_LIGHT = RGBColor(0xB8, 0xA9, 0xD4)
TEAL_BRAND = RGBColor(0x5B, 0xBF, 0xCF)
NAVY = RGBColor(0x1E, 0x1B, 0x4B)
DARK_TEXT = RGBColor(0x2D, 0x2A, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_POS = RGBColor(0x2E, 0xCC, 0x71)
RED_NEG = RGBColor(0xE7, 0x4C, 0x3C)

SCRIPTS_DIR = os.path.dirname(os.path.abspath(__file__))


def add_gradient_bar(slide, prs, position='top'):
    """Add a branded gradient bar to the slide."""
    if position == 'top':
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.15),
            prs.slide_width, Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE_BRAND
    shape.line.fill.background()


def add_logo_text(slide, x, y):
    """Add ModelIt! logo text."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.5), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "ModelIt!"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PURPLE_BRAND


def add_slide_number(slide, num, total=9):
    """Add slide number in bottom-right."""
    box = slide.shapes.add_textbox(
        Inches(8.5), Inches(7.0), Inches(1.2), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = PURPLE_LIGHT
    p.alignment = PP_ALIGN.RIGHT


def add_section_title(slide, title_text, subtitle_text="", prs=None):
    """Add a section title to a slide."""
    if prs:
        add_gradient_bar(slide, prs, 'top')

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.0))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    if subtitle_text:
        sub = tf.add_paragraph()
        sub.text = subtitle_text
        sub.font.size = Pt(16)
        sub.font.color.rgb = DARK_TEXT
        sub.alignment = PP_ALIGN.CENTER


def create_presentation(data):
    """Create a branded 9-slide student presentation from lesson data."""
    lesson_id = data["id"]
    title = data["title"]
    subtitle = data.get("subtitle", "")
    ngss = data.get("ngss", "")
    big_question = data.get("big_question", "")
    objectives = data.get("objectives", [])
    vocabulary = data.get("vocabulary", [])
    components = data.get("components", [])
    relationships = data.get("relationships", [])
    discoveries = data.get("discoveries", [])
    stem_title = data.get("stem_title", "STEM Challenge")
    stem_mission = data.get("stem_mission", "")
    career = data.get("career", "")
    answer = data.get("answer", "")

    # Grade display
    grade_num = lesson_id.split("-")[0].replace("G", "").replace("K", "K")
    if "L2" in grade_num:
        grade_num = grade_num.replace("L2", "")
        level_text = " Level 2"
    else:
        level_text = ""
    if grade_num == "K":
        grade_display = f"Kindergarten{level_text}"
    elif grade_num.isdigit():
        grade_display = f"Grade {int(grade_num)}{level_text}"
    else:
        grade_display = f"Grade {grade_num}{level_text}"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: COVER ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Purple header band
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.5))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = NAVY
    header_bg.line.fill.background()

    # ModelIt! logo
    logo = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.6))
    tf = logo.text_frame
    p = tf.paragraphs[0]
    p.text = "ModelIt!"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEAL_BRAND

    # Doc type
    doc = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(4), Inches(0.4))
    tf = doc.text_frame
    p = tf.paragraphs[0]
    p.text = "Student Presentation"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE_BRAND
    p.alignment = PP_ALIGN.CENTER

    sub = tf.add_paragraph()
    sub.text = subtitle
    sub.font.size = Pt(16)
    sub.font.italic = True
    sub.font.color.rgb = DARK_TEXT
    sub.alignment = PP_ALIGN.CENTER

    # Badges
    badge_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    tf = badge_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{grade_display}  |  NGSS: {ngss}  |  {lesson_id}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PURPLE_BRAND
    p.alignment = PP_ALIGN.CENTER

    add_gradient_bar(slide, prs, 'bottom')
    add_slide_number(slide, 1)

    # ========== SLIDE 2: LEARNING OBJECTIVES ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "What You Will Learn Today", prs=prs)

    # Objectives
    obj_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5), Inches(4.5))
    tf = obj_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Learning Goals"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PURPLE_BRAND

    for obj in objectives[:4]:
        p = tf.add_paragraph()
        p.text = f"  \u2022  {obj}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)

    # Vocabulary
    vocab_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.2), Inches(4.2), Inches(4.5))
    tf = vocab_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Vocabulary"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL_BRAND

    for term, defn in vocabulary[:4]:
        p = tf.add_paragraph()
        p.text = f"  {term}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_before = Pt(6)
        p = tf.add_paragraph()
        p.text = f"    {defn}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT

    add_logo_text(slide, 0.3, 7.0)
    add_slide_number(slide, 2)

    # ========== SLIDE 3: THE BIG QUESTION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bar(slide, prs, 'top')

    q_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(3))
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Big Question"
    p.font.size = Pt(18)
    p.font.color.rgb = TEAL_BRAND
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = big_question
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(20)

    add_logo_text(slide, 0.3, 7.0)
    add_slide_number(slide, 3)

    # ========== SLIDE 4: COMPONENTS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Step 1: LOCATE Your Components", prs=prs)

    ext_comps = [c for c in components if c[2]]
    int_comps = [c for c in components if not c[2]]

    # External column
    ext_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(4.5), Inches(4.5))
    tf = ext_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EXTERNAL (We Control)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL_BRAND

    for comp in ext_comps:
        p = tf.add_paragraph()
        p.text = f"  \u25B6  {comp[0]}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_before = Pt(8)
        p = tf.add_paragraph()
        p.text = f"      {comp[1][:80]}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_TEXT

    # Internal column
    int_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.3), Inches(4.5), Inches(4.5))
    tf = int_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "INTERNAL (System Response)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PURPLE_BRAND

    for comp in int_comps:
        p = tf.add_paragraph()
        p.text = f"  \u25B6  {comp[0]}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_before = Pt(8)
        p = tf.add_paragraph()
        p.text = f"      {comp[1][:80]}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_TEXT

    add_logo_text(slide, 0.3, 7.0)
    add_slide_number(slide, 4)

    # ========== SLIDE 5: RELATIONSHIPS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Step 2: ESTABLISH Relationships", prs=prs)

    rel_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8.5), Inches(4.5))
    tf = rel_box.text_frame
    tf.word_wrap = True

    for i, rel in enumerate(relationships[:5]):
        rel_name = rel[0]
        rel_type = rel[1]
        rel_desc = rel[2] if len(rel) > 2 else ""

        is_positive = "+" in rel_type or "POSITIVE" in rel_type.upper()
        symbol = "+" if is_positive else "\u2212"
        color = GREEN_POS if is_positive else RED_NEG

        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"  {symbol}  {rel_name}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_before = Pt(10)

        p = tf.add_paragraph()
        p.text = f"       {rel_desc[:100]}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT

    add_logo_text(slide, 0.3, 7.0)
    add_slide_number(slide, 5)

    # ========== SLIDE 6: SIMULATION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Step 3: VISUALIZE & EVALUATE", "Run the simulation and observe patterns", prs=prs)

    scenarios = data.get("sim_scenarios", data.get("scenarios", []))
    sc_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.5), Inches(4))
    tf = sc_box.text_frame
    tf.word_wrap = True

    for i, sc in enumerate(scenarios[:3]):
        sc_title = sc[0] if isinstance(sc, (list, tuple)) else str(sc)
        sc_setup = sc[1] if isinstance(sc, (list, tuple)) and len(sc) > 1 else ""

        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"Scenario {i+1}: {sc_title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEAL_BRAND
        p.space_before = Pt(14)

        p = tf.add_paragraph()
        p.text = f"    {sc_setup[:120]}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT

    add_logo_text(slide, 0.3, 7.0)
    add_slide_number(slide, 6)

    # ========== SLIDE 7: DISCOVERIES ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "What Did We Discover?", prs=prs)

    disc_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8.5), Inches(4.5))
    tf = disc_box.text_frame
    tf.word_wrap = True

    for i, disc in enumerate(discoveries[:5]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"  \u2713  {disc}"
        p.font.size = Pt(16)
        p.font.color.rgb = NAVY
        p.space_before = Pt(12)

    add_logo_text(slide, 0.3, 7.0)
    add_slide_number(slide, 7)

    # ========== SLIDE 8: STEM CHALLENGE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, f"Step 4: {stem_title}", prs=prs)

    stem_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8.5), Inches(4.5))
    tf = stem_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "YOUR ENGINEERING MISSION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEAL_BRAND

    p = tf.add_paragraph()
    p.text = stem_mission[:250]
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    p.space_before = Pt(10)

    # Career connection
    if career:
        p = tf.add_paragraph()
        p.space_before = Pt(20)
        p.text = "Career Connection"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PURPLE_BRAND

        p = tf.add_paragraph()
        p.text = career[:200]
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT

    add_logo_text(slide, 0.3, 7.0)
    add_slide_number(slide, 8)

    # ========== SLIDE 9: EXIT TICKET ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full navy background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Logo
    logo = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(4), Inches(0.8))
    tf = logo.text_frame
    p = tf.paragraphs[0]
    p.text = "ModelIt!"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEAL_BRAND
    p.alignment = PP_ALIGN.CENTER

    exit_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = exit_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Exit Ticket"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    exit_questions = [
        f"What is the relationship between {components[0][0] if components else 'the external component'} and {components[-1][0] if components else 'the internal component'}?",
        f"What happened in your simulation when you changed {components[0][0] if components else 'the input'}?",
        "What is ONE thing you would add to your model and why?"
    ]

    for i, eq in enumerate(exit_questions, 1):
        p = tf.add_paragraph()
        p.text = f"\n{i}. {eq}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(14)

    # Footer
    footer = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{lesson_id}  |  {grade_display}  |  NGSS: {ngss}  |  Discovery Collective \u00A9 2026"
    p.font.size = Pt(10)
    p.font.color.rgb = PURPLE_LIGHT
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 9)

    return prs


# ============================================================
# Load and generate
# ============================================================

def load_lesson_data(grade, lesson_num):
    """Load lesson data — reuses logic from create_branded_materials.py."""
    grade_lower = grade.lower()
    data_file = None

    # Map grade to data file
    grade_map = {
        "k": "lesson_data_GK.py", "0k": "lesson_data_GK.py",
        "01": "lesson_data_G01.py", "1": "lesson_data_G01.py",
        "02": "lesson_data_G02.py", "2": "lesson_data_G02.py",
        "03": "lesson_data_G03.py", "3": "lesson_data_G03.py",
        "04": "lesson_data_G04.py", "4": "lesson_data_G04.py",
        "04_l2": "lesson_data_G04_L2.py",
        "05_l2": "lesson_data_G05_L2.py",
        "06": "lesson_data_G06.py", "6": "lesson_data_G06.py",
        "06_l2": "lesson_data_G06_L2.py",
        "07": "lesson_data_G07.py", "7": "lesson_data_G07.py",
        "07_l2": "lesson_data_G07_L2.py",
        "08": "lesson_data_G08.py", "8": "lesson_data_G08.py",
        "08_l2": "lesson_data_G08_L2.py",
    }

    if grade_lower in grade_map:
        data_file = os.path.join(SCRIPTS_DIR, grade_map[grade_lower])
    elif "_" in grade_lower:
        data_file = os.path.join(SCRIPTS_DIR, f"lesson_data_G{grade_lower}.py")

    if not data_file or not os.path.exists(data_file):
        print(f"ERROR: Data file not found for grade '{grade}'")
        return None

    spec = importlib.util.spec_from_file_location(f"ld_{grade}", data_file)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)

    lesson_key = f"L{lesson_num:02d}"
    if hasattr(mod, lesson_key):
        return getattr(mod, lesson_key)

    print(f"WARNING: Lesson {lesson_key} not found in {data_file}")
    return None


def get_output_dir(grade, lesson_num):
    """Get output directory for presentations."""
    base = os.path.join(os.path.dirname(SCRIPTS_DIR), "materials")
    grade_lower = grade.lower()

    if "_" in grade_lower:
        grade_num, level = grade_lower.split("_")
        level_num = level.replace("l", "").replace("L", "")
        grade_folder = f"grade-{grade_num}"
        level_folder = f"level-{level_num}"
        lesson_prefix = f"G{grade_num}L{level_num}"
        lesson_folder = f"{lesson_prefix}-L{lesson_num:02d}"
        output_dir = os.path.join(base, grade_folder, level_folder, lesson_folder)
    elif grade_lower in ["k", "0k"]:
        lesson_folder = f"GK-L{lesson_num:02d}"
        output_dir = os.path.join(base, "grade-K", lesson_folder)
    elif len(grade_lower) <= 2:
        g = grade_lower.zfill(2)
        lesson_folder = f"G{g}-L{lesson_num:02d}"
        output_dir = os.path.join(base, f"grade-{g}", lesson_folder)
    else:
        lesson_folder = f"G{grade_lower}-L{lesson_num:02d}"
        output_dir = os.path.join(base, f"grade-{grade_lower}", lesson_folder)

    os.makedirs(output_dir, exist_ok=True)
    return output_dir, lesson_folder


def main():
    parser = argparse.ArgumentParser(description="ModelIt! Student Presentation Generator")
    parser.add_argument("--grade", required=True, help="Grade level (04_L2, 06, etc.)")
    parser.add_argument("--lesson", type=int, help="Lesson number")
    parser.add_argument("--all", action="store_true", help="Generate all lessons")

    args = parser.parse_args()

    if args.all:
        if args.grade.lower().endswith("_l2") and args.grade[:2].isdigit() and int(args.grade[:2]) <= 8:
            lessons = range(1, 8)
        elif args.grade.lower() in ["ne", "natures-engineers"]:
            lessons = range(1, 13)
        else:
            lessons = range(1, 11)
    elif args.lesson:
        lessons = [args.lesson]
    else:
        print("ERROR: Specify --lesson N or --all")
        sys.exit(1)

    for lesson_num in lessons:
        print(f"\nGenerating presentation: Grade {args.grade}, Lesson {lesson_num}")

        data = load_lesson_data(args.grade, lesson_num)
        if not data:
            print(f"  SKIP: No data for lesson {lesson_num}")
            continue

        output_dir, lesson_folder = get_output_dir(args.grade, lesson_num)

        prs = create_presentation(data)
        filepath = os.path.join(output_dir, f"{lesson_folder}-Student-Presentation.pptx")
        prs.save(filepath)
        print(f"  [OK] {filepath}")

    print("\nDone!")


if __name__ == "__main__":
    main()
